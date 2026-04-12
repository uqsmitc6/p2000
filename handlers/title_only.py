"""
Title Only Handler — Slides with just a title and possibly images/diagrams.

Template layout 43 "Title Only":
    idx  0 — Title
    idx 17 — Footer (programme name)
    idx 18 — Slide number

Detection: slide has a title but very little or no body text. The slide
may have images, diagrams, or charts that fill the content area. These
slides are distinguished from Section Dividers by not having the sparse,
structured section-number + title pattern.

Images are extracted from the source slide and placed as free-floating
shapes in the content area below the title.

CRITICAL RULE: Never set font properties on placeholders.
"""

import io
import re
from pptx.util import Emu
from handlers.base import SlideHandler
from utils.extractor import extract_text_elements, extract_images, extract_shapes_with_text


class TitleOnlyHandler(SlideHandler):

    name = "Title Only"
    description = "Slide with just a title (and possibly images that need manual placement)"
    layout_name = "Title Only"
    layout_index = 43

    PH_TITLE = 0
    PH_FOOTER = 17
    PH_SLIDE_NUM = 18

    PLACEHOLDER_NOISE = [
        "[click", "[add", "[your", "[insert", "[subtitle",
        "click to add", "click icon", "click to edit",
        "\u00a9 the university", "this content is protected",
        "cricos code",
    ]

    # --- Detection ---

    def detect(self, slide, slide_index: int) -> float:
        """
        Detect slides that have a title but minimal body text.
        These are typically diagram/image slides where the visual
        content can't be auto-extracted.
        """
        if slide_index == 0:
            return 0.0

        # Check source layout name
        layout_name = slide.slide_layout.name.lower()
        if "title only" in layout_name:
            return 0.70

        shapes = extract_shapes_with_text(slide)
        images = extract_images(slide)

        if not shapes:
            # No text at all — might be a blank or image-only slide
            if images:
                return 0.50  # Has images but no text — Title Only is best fit
            return 0.0

        # Filter noise
        meaningful = [
            s for s in shapes
            if s["text"].strip()
            and not any(p in s["text"].lower() for p in self.PLACEHOLDER_NOISE)
            and len(s["text"].strip()) > 1
            and not re.match(r"^\d{1,3}$", s["text"].strip())
        ]

        if not meaningful:
            return 0.0

        # Separate title from body: title is likely the topmost short text
        sorted_by_top = sorted(meaningful, key=lambda s: s["top"] or 0)
        title_candidate = None
        body_text_total = 0

        for s in sorted_by_top:
            if title_candidate is None and len(s["text"]) < 120:
                title_candidate = s
            else:
                body_text_total += len(s["text"])

        has_images = bool(images)

        # Title only: has a title but very little body text
        if title_candidate and body_text_total < 30:
            if has_images:
                return 0.68  # Title + images, minimal text → strong Title Only
            elif len(meaningful) == 1:
                return 0.45  # Just a title, no images
            return 0.0

        # Also catch slides with 2 short shapes (title + subtitle-like)
        if len(meaningful) == 2 and has_images:
            lengths = sorted(len(s["text"]) for s in meaningful)
            if lengths[0] < 50 and lengths[1] < 120:
                return 0.55

        return 0.0

    # --- Extraction ---

    def extract_content(self, slide, slide_index: int) -> dict:
        """
        Extract the title and any images from the slide.
        Images are stored as blobs with position/size info for placement.
        """
        shapes = extract_shapes_with_text(slide)

        result = {
            "title": "",
            "footer": "",
            "has_images": False,
            "images": [],  # list of {"blob": bytes, "width": emu, "height": emu}
        }

        # Extract images with their data
        images = extract_images(slide)
        result["has_images"] = bool(images)
        for img in images:
            if img.get("blob"):
                result["images"].append({
                    "blob": img["blob"],
                    "width": img.get("width"),
                    "height": img.get("height"),
                })

        if not shapes:
            return result

        # Filter noise
        filtered = []
        for s in shapes:
            text = s["text"].strip()
            if not text:
                continue
            if any(p in text.lower() for p in self.PLACEHOLDER_NOISE):
                continue
            if re.match(r"^\d{1,3}$", text) and len(text) <= 3:
                continue
            filtered.append(s)

        # Find title
        title_shape = None
        footer_text = ""

        for s in filtered:
            if s["is_placeholder"]:
                idx = s["shape_idx"]
                if idx in (0, 3, 15):
                    title_shape = s
                elif idx in (17,):
                    if not footer_text:
                        footer_text = s["text"]
                elif idx in (18,):
                    pass
            elif not title_shape:
                # First non-footer text becomes title
                if not self._is_footer_text(s["text"]) and len(s["text"]) < 120:
                    title_shape = s

        result["footer"] = footer_text

        if title_shape:
            result["title"] = title_shape["text"]
        elif filtered:
            # Fallback: largest font
            with_font = [s for s in filtered if s["font_size"]]
            if with_font:
                result["title"] = max(with_font, key=lambda s: s["font_size"])["text"]

        return result

    def _is_footer_text(self, text: str) -> bool:
        footer_patterns = [
            r"(?i)cricos", r"(?i)hbis\s+innovation",
            r"(?i)executive\s+education",
            r"(?i)presentation\s+title",
        ]
        return any(re.search(p, text) for p in footer_patterns)

    # --- Output ---

    # Content area dimensions (below title, above footer)
    # Title bottom ~1.49in, footer top ~6.98in, left margin ~0.52in
    CONTENT_TOP = Emu(1500000)      # ~1.64in — a little below title bottom
    CONTENT_LEFT = Emu(479424)      # matches title left margin
    CONTENT_WIDTH = Emu(11233151)   # slide width minus margins
    CONTENT_HEIGHT = Emu(4800000)   # ~5.25in available height
    FOOTER_TOP = Emu(6381750)       # footer starts here

    def fill_slide(self, slide, content: dict) -> None:
        """
        Fill Title Only placeholder and place extracted images.
        NEVER set font properties.
        """
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        if content.get("title") and self.PH_TITLE in placeholders:
            placeholders[self.PH_TITLE].text = content["title"]

        # Place extracted images in the content area
        images = content.get("images", [])
        if images:
            self._place_images(slide, images)

    def _place_images(self, slide, images: list) -> None:
        """
        Place extracted images in the content area below the title.

        For a single image: centre it in the available space, scaling
        to fit while maintaining aspect ratio.
        For multiple images: arrange side by side if they fit, otherwise
        just place the largest one.
        """
        if not images:
            return

        avail_w = self.CONTENT_WIDTH
        avail_h = self.CONTENT_HEIGHT

        if len(images) == 1:
            self._place_single_image(slide, images[0], avail_w, avail_h)
        else:
            # Multiple images — place largest one centred, skip the rest
            # (most slides have one main image; extras are usually logos)
            largest = max(images, key=lambda i: (i.get("width") or 0) * (i.get("height") or 0))
            self._place_single_image(slide, largest, avail_w, avail_h)

    def _place_single_image(self, slide, img_data: dict, avail_w, avail_h) -> None:
        """Place a single image centred in the content area, scaled to fit."""
        from PIL import Image as PILImage

        blob = img_data["blob"]
        image_stream = io.BytesIO(blob)

        # Get actual pixel dimensions to compute aspect ratio
        try:
            pil_img = PILImage.open(io.BytesIO(blob))
            px_w, px_h = pil_img.size
        except Exception:
            # Fallback: use EMU dimensions from extraction, or assume 4:3
            px_w = img_data.get("width") or 4
            px_h = img_data.get("height") or 3

        if px_w == 0 or px_h == 0:
            return

        aspect = px_w / px_h

        # Scale to fit available area maintaining aspect ratio
        if aspect >= (avail_w / avail_h):
            # Width-constrained
            img_w = avail_w
            img_h = int(avail_w / aspect)
        else:
            # Height-constrained
            img_h = avail_h
            img_w = int(avail_h * aspect)

        # Centre in content area
        left = self.CONTENT_LEFT + (avail_w - img_w) // 2
        top = self.CONTENT_TOP + (avail_h - img_h) // 2

        slide.shapes.add_picture(image_stream, left, top, img_w, img_h)

    def get_placeholder_map(self) -> dict:
        return {
            "title": self.PH_TITLE,
            "footer": self.PH_FOOTER,
            "slide_num": self.PH_SLIDE_NUM,
        }
