"""
Cover 1 Handler — White background, text on left, UQ graphic device.

Placeholders:
    idx 0  — Title (48pt, inherited from theme)
    idx 10 — Subtitle / description (18pt bold, inherited)
    idx 11 — Entity name / school (16pt, inherited)

Detection: Cover slides are typically the first slide and contain a
programme title, entity name, and optionally presenter/date info.
"""

import re
from handlers.base import SlideHandler
from utils.extractor import extract_text_elements, extract_images


class Cover1Handler(SlideHandler):

    name = "Cover 1"
    description = "White background cover slide with text on left"
    layout_name = "Cover 1"
    layout_index = 0

    # Placeholder indices in the template layout
    PH_TITLE = 0
    PH_SUBTITLE = 10
    PH_ENTITY = 11

    # --- Detection ---

    PLACEHOLDER_NOISE = [
        "[program name]", "[insert", "[click", "[add", "[your",
        "click to add", "click icon", "\u00a9 the university",
        "this content is protected", "may not be shared",
    ]

    PRESENTER_PATTERNS = [
        r"(?i)^(a/?)?professor\s",
        r"(?i)^(associate|assistant)\s+professor\s",
        r"(?i)^dr\.?\s",
        r"(?i)^presenter[s]?[:\s]",
        r"(?i)^professors?\s.*\sand\s",
    ]

    ENTITY_PATTERNS = [
        r"(?i)uq\s", r"(?i)university\s+of\s+queensland",
        r"(?i)business\s+school", r"(?i)executive\s+education",
        r"(?i)faculty\s+of", r"(?i)school\s+of",
        r"(?i)institute\s+of", r"(?i)centre\s+for",
        r"(?i)open\s+program",
    ]

    DATE_PATTERNS = [
        r"(?i)^(january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{4}$",
        r"(?i)^day\s+\d",
        r"(?i)^\d{1,2}\s+(january|february|march|april|may|june|july|august|september|october|november|december)",
        r"(?i)^(semester|week|session|module)\s+\d",
        r"(?i)^\d{4}$",
    ]

    def detect(self, slide, slide_index: int) -> float:
        """
        Heuristic: first slide with a large-font title and no dominant
        image is very likely a cover slide.
        """
        # Cover is the FIRST slide of the deck. There's only ever one.
        if slide_index == 0:
            return 0.9

        # Non-first slides are almost never covers.
        return 0.0

    # --- Extraction & Classification ---

    def extract_content(self, slide, slide_index: int) -> dict:
        """Extract and classify content from an input cover slide."""
        texts = extract_text_elements(slide)
        images = extract_images(slide)

        # Expand multi-line text blobs into individual lines
        expanded = []
        for t in texts:
            lines = t["text"].split("\n") if "\n" in t["text"] else [t["text"]]
            for line in lines:
                line = line.strip()
                if line:
                    expanded.append({**t, "text": line})

        # Filter placeholder/template noise
        filtered = [
            t for t in expanded
            if not any(p in t["text"].lower().strip() for p in self.PLACEHOLDER_NOISE)
            and len(t["text"].strip()) > 1
        ]

        # Classify each text fragment
        presenters = []
        entities = []
        dates = []
        title_candidates = []

        for t in filtered:
            text = t["text"]

            if any(re.search(p, text) for p in self.PRESENTER_PATTERNS):
                presenters.append(t)
            elif any(re.search(p, text) for p in self.DATE_PATTERNS):
                dates.append(t)
            elif any(re.search(p, text) for p in self.ENTITY_PATTERNS):
                entities.append(t)
            else:
                title_candidates.append(t)

        # Largest font → title; next → subtitle
        title_candidates.sort(key=lambda t: (t["font_size"] or 0), reverse=True)

        result = {
            "title": title_candidates[0]["text"] if title_candidates else "",
            "subtitle": title_candidates[1]["text"] if len(title_candidates) >= 2 else "",
            "entity": entities[0]["text"] if entities else "",
            "presenter": "",
            "date_info": "",
            "extra_texts": [t["text"] for t in title_candidates[2:]] if len(title_candidates) > 2 else [],
            "has_image": len(images) > 0,
        }

        # Deduplicate presenters
        if presenters:
            seen = set()
            unique = []
            for p in presenters:
                clean = re.sub(r"(?i)^presenters?[:\s]*", "", p["text"]).strip()
                if clean not in seen:
                    seen.add(clean)
                    unique.append(clean)
            result["presenter"] = "\n".join(unique)

        # Date info
        if dates:
            result["date_info"] = " | ".join(d["text"] for d in dates)

        # Build subtitle from presenter + date if no explicit subtitle
        if not result["subtitle"]:
            parts = []
            if result["presenter"]:
                parts.append(result["presenter"])
            if result["date_info"]:
                parts.append(result["date_info"])
            result["subtitle"] = "\n".join(parts)

        # If title is very long, try to split at a natural break point
        # (colon, dash, em-dash) so the first part stays in the big title
        # and the rest becomes the first line of the subtitle.
        if result["title"]:
            title_lines = self._estimate_title_lines(result["title"])
            if title_lines > self.MAX_TITLE_LINES:
                m = self.TITLE_SPLIT_PATTERN.match(result["title"])
                if m:
                    result["title"] = m.group(1).rstrip() + ":"
                    overflow = m.group(2).strip()
                    # Prepend overflow to subtitle
                    if result["subtitle"]:
                        result["subtitle"] = overflow + "\n" + result["subtitle"]
                    else:
                        result["subtitle"] = overflow

        return result

    # --- Title splitting for long titles ---

    # If a title would wrap to more than this many lines at 48pt, try to
    # split it at a natural break point (colon, dash, em-dash) so the first
    # part stays in the title placeholder and the rest becomes the first
    # line of the subtitle placeholder (displayed at 18pt bold).
    MAX_TITLE_LINES = 3

    TITLE_SPLIT_PATTERN = re.compile(
        r'^(.{15,}?)\s*[:–—-]\s*(.+)$', re.DOTALL
    )

    # --- Title line estimation ---

    # Calibrated against real outputs: at 48pt Arial in a 5.10" placeholder,
    # approximately 20 characters fit per line with word-wrap.
    CHARS_PER_LINE = 20

    # How much to shift the subtitle down per extra title line beyond 2.
    # 0.5 inches per line keeps the visual gap consistent without pushing
    # the subtitle off the slide. (0.5" ≈ 457200 EMU)
    SHIFT_PER_EXTRA_LINE_EMU = 457200

    def _estimate_title_lines(self, title: str) -> int:
        """
        Estimate how many lines a title will wrap to in the placeholder.
        Uses word-boundary wrapping for accuracy.
        """
        if not title:
            return 0

        words = title.split()
        lines = 1
        current_line_len = 0

        for word in words:
            word_len = len(word)
            if current_line_len == 0:
                current_line_len = word_len
            elif current_line_len + 1 + word_len <= self.CHARS_PER_LINE:
                current_line_len += 1 + word_len  # space + word
            else:
                lines += 1
                current_line_len = word_len

        return lines

    # --- Output ---

    def fill_slide(self, slide, content: dict) -> None:
        """
        Fill Cover 1 placeholders with extracted content.
        ONLY sets .text — all formatting inherits from the template.

        Dynamically adjusts subtitle position when the title is long
        enough to overflow its placeholder (3+ lines at 48pt).
        """
        from pptx.util import Emu

        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Title
        title_lines = 0
        if content.get("title") and self.PH_TITLE in placeholders:
            placeholders[self.PH_TITLE].text = content["title"]
            title_lines = self._estimate_title_lines(content["title"])

        # Calculate subtitle position adjustment for long titles.
        # 2 lines = baseline (no shift). Each extra line beyond 2
        # pushes the subtitle down by SHIFT_PER_EXTRA_LINE_EMU.
        extra_lines = max(0, title_lines - 2)
        subtitle_shift = extra_lines * self.SHIFT_PER_EXTRA_LINE_EMU

        # Subtitle
        if content.get("subtitle") and self.PH_SUBTITLE in placeholders:
            ph = placeholders[self.PH_SUBTITLE]

            # Adjust position if title overflows.
            # IMPORTANT: placeholders inherit geometry from the layout.
            # When we modify .top, python-pptx creates an explicit xfrm
            # element but may not preserve inherited left/width/height.
            # We must capture and re-set all geometry values.
            if subtitle_shift > 0:
                orig_left = ph.left
                orig_width = ph.width
                orig_height = ph.height
                ph.top = ph.top + subtitle_shift
                ph.left = orig_left
                ph.width = orig_width
                ph.height = orig_height
                # No leading newline needed — position handles the gap
                ph.text = content["subtitle"]
            else:
                # Standard case: prepend blank line for visual spacing
                ph.text = "\n" + content["subtitle"]

        # Entity
        if content.get("entity") and self.PH_ENTITY in placeholders:
            placeholders[self.PH_ENTITY].text = content["entity"]

    def get_placeholder_map(self) -> dict:
        return {
            "title": self.PH_TITLE,
            "subtitle": self.PH_SUBTITLE,
            "entity": self.PH_ENTITY,
        }
