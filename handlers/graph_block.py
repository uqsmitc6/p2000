"""
Graph with Block Handler — Graph/chart with coloured text panel.

Template layouts (three colour variants, identical structure):
    idx 29 — Graph with Dark Purple Block
    idx 30 — Graph with Neutral Block
    idx 31 — Graph with Grey Block

Placeholder mapping (same across all three variants):
    idx  0 — Title (left side, top)
    idx 31 — Subtitle (left side, below title)
    idx 10 — Content/Graph area (left side, below subtitle, 6.14in × 3.36in)
    idx 14 — Block text (right side, full-height coloured panel, 5.85in × 5.59in)
    idx 19 — Footer
    idx 20 — Slide number

Detection: slide has a graph/chart/image on one side AND a block of
descriptive text on the other. Distinguished from Text with Image by
the presence of a substantial text block alongside the visual content.

Colour variant selection defaults to Dark Purple (most common in UQ
branding) unless the source slide name hints at a different variant.

CRITICAL RULE: Never set font properties on placeholders.
"""

import io
import re
from pptx.util import Emu
from handlers.base import SlideHandler
from utils.extractor import extract_shapes_with_text, extract_images


# Colour variant configs: (layout_index, variant_name)
VARIANT_DARK_PURPLE = (29, "dark_purple")
VARIANT_NEUTRAL = (30, "neutral")
VARIANT_GREY = (31, "grey")

# Default to Dark Purple (primary UQ branding)
DEFAULT_VARIANT = VARIANT_DARK_PURPLE


class GraphBlockHandler(SlideHandler):

    name = "Graph with Block"
    description = "Graph or chart with coloured text block panel"
    layout_name = "Graph with Dark Purple Block"
    layout_index = 29  # Default, overridden by get_layout_index()

    PH_TITLE = 0
    PH_SUBTITLE = 31
    PH_CONTENT = 10     # Graph/chart area (left)
    PH_BLOCK_TEXT = 14   # Coloured text panel (right)
    PH_FOOTER = 19
    PH_SLIDE_NUM = 20

    PLACEHOLDER_NOISE = [
        "[click", "[add", "[your", "[insert", "[subtitle",
        "click to add", "click icon", "click to edit",
        "\u00a9 the university", "this content is protected",
        "cricos code",
    ]

    FOOTER_PATTERNS = [
        r"(?i)cricos", r"(?i)hbis\s+innovation",
        r"(?i)executive\s+education",
        r"(?i)presentation\s+title",
    ]

    SLIDE_WIDTH = 12192000
    SLIDE_MID = SLIDE_WIDTH // 2

    # Content area dimensions for image placement
    CONTENT_LEFT = Emu(479424)       # 0.52in
    CONTENT_TOP = Emu(3310128)       # 3.62in
    CONTENT_WIDTH = Emu(5616576)     # 6.14in
    CONTENT_HEIGHT = Emu(3073344)    # 3.36in

    # --- Detection ---

    def detect(self, slide, slide_index: int) -> float:
        """
        Detect slides suitable for graph-with-block layout.

        Triggers:
        - Source layout name contains "graph" and "block"
        - Slide has an image/chart on one side + substantial text on the other
        """
        if slide_index == 0:
            return 0.0

        # Check source layout name
        layout_name = slide.slide_layout.name.lower()
        if "graph" in layout_name and "block" in layout_name:
            return 0.75

        # Heuristic: image + substantial separate text block
        images = extract_images(slide)
        shapes = extract_shapes_with_text(slide)
        meaningful = self._get_meaningful_shapes(shapes)

        if not images or not meaningful:
            return 0.0

        # Need at least one real content image
        real_images = [
            img for img in images
            if img.get("width") and img.get("height")
            and img["width"] > self.SLIDE_WIDTH * 0.15
            and img["height"] > 6858000 * 0.15
        ]

        if not real_images:
            return 0.0

        # Separate title from body
        title_shape = self._find_title(meaningful)
        body_shapes = [s for s in meaningful if s is not title_shape]

        if not body_shapes:
            return 0.0

        # Look for a substantial text block (>100 chars) — the block panel text
        block_candidates = [s for s in body_shapes if len(s["text"]) >= 100]
        if not block_candidates:
            return 0.0

        # The block text should be positioned differently from the graph/image
        # (i.e. one on left, one on right)
        for block in block_candidates:
            block_centre = (block["left"] or 0) + (block["width"] or 0) / 2
            for img in real_images:
                img_centre = (img["left"] or 0) + (img["width"] or 0) / 2
                # They should be on opposite sides of the slide
                if abs(block_centre - img_centre) > self.SLIDE_WIDTH * 0.25:
                    return 0.50

        return 0.0

    # --- Extraction ---

    def extract_content(self, slide, slide_index: int) -> dict:
        """
        Extract title, subtitle, block text, and any graph/image content.
        Also determines colour variant from source layout name.
        """
        shapes = extract_shapes_with_text(slide)
        images = extract_images(slide)

        result = {
            "title": "",
            "subtitle": "",
            "content": "",          # Left-side content (text below subtitle)
            "block_text": "",       # Right-side coloured panel text
            "footer": "",
            "has_images": bool(images),
            "images": [],
            "_variant": DEFAULT_VARIANT,
        }

        # Determine colour variant from source layout name
        layout_name = slide.slide_layout.name.lower()
        if "neutral" in layout_name:
            result["_variant"] = VARIANT_NEUTRAL
        elif "grey" in layout_name or "gray" in layout_name:
            result["_variant"] = VARIANT_GREY
        else:
            result["_variant"] = VARIANT_DARK_PURPLE

        # Extract images
        for img in images:
            if img.get("blob"):
                result["images"].append({
                    "blob": img["blob"],
                    "width": img.get("width"),
                    "height": img.get("height"),
                })

        if not shapes:
            return result

        meaningful = self._get_meaningful_shapes(shapes)
        if not meaningful:
            return result

        # Separate by role
        title_shape = None
        body_shapes = []
        footer_text = ""

        for s in meaningful:
            if s["is_placeholder"]:
                idx = s["shape_idx"]
                if idx in (0, 3, 15):
                    title_shape = s
                elif idx in (17, 19):
                    if not footer_text:
                        footer_text = s["text"]
                elif idx in (18, 20):
                    pass
                else:
                    body_shapes.append(s)
            elif self._is_footer_text(s["text"]):
                if not footer_text:
                    footer_text = s["text"]
            else:
                body_shapes.append(s)

        result["footer"] = footer_text

        # Fallback title
        if not title_shape and body_shapes:
            with_font = [s for s in body_shapes if s["font_size"]]
            if with_font:
                candidate = max(with_font, key=lambda s: s["font_size"])
                if len(candidate["text"]) < 120:
                    title_shape = candidate
                    body_shapes = [s for s in body_shapes if s is not title_shape]

        if title_shape:
            result["title"] = title_shape["text"]

        # Subtitle detection
        body_shapes.sort(key=lambda s: (s["top"] or 0, s["left"] or 0))
        if body_shapes and title_shape:
            candidate = body_shapes[0]
            is_short = len(candidate["text"]) < 100
            is_near_top = candidate["top"] < (title_shape["top"] + title_shape["height"] * 2)
            has_more = len(body_shapes) > 1
            if is_short and is_near_top and has_more:
                result["subtitle"] = candidate["text"]
                body_shapes = body_shapes[1:]

        # Split remaining body into left (content) and right (block text)
        # based on horizontal position
        left_texts = []
        right_texts = []

        for s in body_shapes:
            shape_centre = (s["left"] or 0) + (s["width"] or 0) / 2
            if shape_centre < self.SLIDE_MID:
                left_texts.append(s["text"])
            else:
                right_texts.append(s["text"])

        result["content"] = "\n".join(left_texts)
        result["block_text"] = "\n".join(right_texts)

        # If all text ended up on one side, use the longest block as block_text
        if result["block_text"] and not result["content"]:
            # All on right — keep as is
            pass
        elif result["content"] and not result["block_text"]:
            # All on left — assign longest text group to block
            all_shapes = sorted(body_shapes, key=lambda s: len(s["text"]), reverse=True)
            if len(all_shapes) >= 2:
                result["block_text"] = all_shapes[0]["text"]
                result["content"] = "\n".join(s["text"] for s in all_shapes[1:])
            elif all_shapes:
                result["block_text"] = all_shapes[0]["text"]
                result["content"] = ""

        return result

    # --- Dynamic layout ---

    def get_layout_index(self, content: dict) -> int:
        """Return the appropriate colour variant layout index."""
        variant = content.get("_variant", DEFAULT_VARIANT)
        return variant[0]

    # --- Output ---

    def fill_slide(self, slide, content: dict) -> None:
        """
        Fill graph-with-block placeholders. Images are placed as
        free-floating shapes in the content area (left side).
        NEVER set font properties.
        """
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        if content.get("title") and self.PH_TITLE in placeholders:
            placeholders[self.PH_TITLE].text = content["title"]

        if content.get("subtitle") and self.PH_SUBTITLE in placeholders:
            placeholders[self.PH_SUBTITLE].text = content["subtitle"]

        if content.get("content") and self.PH_CONTENT in placeholders:
            placeholders[self.PH_CONTENT].text = content["content"]

        if content.get("block_text") and self.PH_BLOCK_TEXT in placeholders:
            placeholders[self.PH_BLOCK_TEXT].text = content["block_text"]

        # Place images in the content area (left side, below subtitle)
        images = content.get("images", [])
        if images:
            self._place_images(slide, images)

    def _place_images(self, slide, images: list) -> None:
        """Place extracted images in the left content area."""
        from PIL import Image as PILImage

        if not images:
            return

        # Place the largest image, scaled to fit content area
        largest = max(images, key=lambda i: (i.get("width") or 0) * (i.get("height") or 0))
        blob = largest.get("blob")
        if not blob:
            return

        try:
            pil_img = PILImage.open(io.BytesIO(blob))
            px_w, px_h = pil_img.size
        except Exception:
            px_w = largest.get("width") or 4
            px_h = largest.get("height") or 3

        if px_w == 0 or px_h == 0:
            return

        aspect = px_w / px_h
        avail_w = self.CONTENT_WIDTH
        avail_h = self.CONTENT_HEIGHT

        if aspect >= (avail_w / avail_h):
            img_w = avail_w
            img_h = int(avail_w / aspect)
        else:
            img_h = avail_h
            img_w = int(avail_h * aspect)

        left = self.CONTENT_LEFT + (avail_w - img_w) // 2
        top = self.CONTENT_TOP + (avail_h - img_h) // 2

        image_stream = io.BytesIO(blob)
        slide.shapes.add_picture(image_stream, left, top, img_w, img_h)

    # --- Helpers ---

    def _get_meaningful_shapes(self, shapes: list) -> list:
        meaningful = []
        for s in shapes:
            text = s["text"].strip()
            if not text:
                continue
            if any(p in text.lower() for p in self.PLACEHOLDER_NOISE):
                continue
            if len(text) <= 1:
                continue
            if re.match(r"^\d{1,3}$", text) and len(text) <= 3:
                continue
            meaningful.append(s)
        return meaningful

    def _find_title(self, shapes: list):
        for s in shapes:
            if s["is_placeholder"] and s["shape_idx"] in (0, 3, 15):
                return s
        sorted_by_top = sorted(shapes, key=lambda s: s["top"] or 0)
        for s in sorted_by_top:
            if len(s["text"]) < 120:
                if s.get("font_size") and s["font_size"] >= 20:
                    return s
                elif not s.get("font_size") and len(s["text"]) < 80:
                    return s
        return None

    def _is_footer_text(self, text: str) -> bool:
        return any(re.search(p, text) for p in self.FOOTER_PATTERNS)

    def get_placeholder_map(self) -> dict:
        return {
            "title": self.PH_TITLE,
            "subtitle": self.PH_SUBTITLE,
            "content": self.PH_CONTENT,
            "block_text": self.PH_BLOCK_TEXT,
            "footer": self.PH_FOOTER,
            "slide_num": self.PH_SLIDE_NUM,
        }
