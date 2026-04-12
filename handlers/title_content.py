"""
Title and Content Handler — The standard content slide.

Template placeholders (layout index 6, "Title and Content"):
    idx  0 — Title
    idx 31 — Subtitle (optional — if absent, content shifts up)
    idx 10 — Content body (full width, 12.28" × 4.51")
    idx 17 — Footer (programme name)
    idx 18 — Slide number

Detection: This is the default/fallback handler. Any slide that isn't
a cover, section divider, thank you, or contents page is likely a
content slide. Most slides in any deck will be this type.

Design concerns (future):
    - Wall-of-text risk: content body is huge (12.28" × 4.51").
      Future versions should enforce visual communication rules
      (max bullet points, text density limits, suggest splitting).
    - Subtitle is optional: when absent, content shifts up 0.55"
      to reclaim the dead space.
"""

import re
from handlers.base import SlideHandler
from utils.extractor import extract_text_elements, extract_images, extract_shapes_with_text, extract_rich_paragraphs


class TitleContentHandler(SlideHandler):

    name = "Title and Content"
    description = "Standard content slide with title, optional subtitle, and body"
    layout_name = "Title and Content"
    layout_index = 6

    # Placeholder indices
    PH_TITLE = 0
    PH_SUBTITLE = 31
    PH_CONTENT = 10
    PH_FOOTER = 17
    PH_SLIDE_NUM = 18

    # Subtitle placeholder height — used to shift content up when no subtitle
    SUBTITLE_HEIGHT_EMU = 502920  # ~0.55 inches

    # Noise to filter
    PLACEHOLDER_NOISE = [
        "[click", "[add", "[your", "[insert", "[subtitle",
        "click to add", "click icon", "click to edit",
        "\u00a9 the university", "this content is protected",
        "cricos code",
    ]

    # --- Detection ---

    def detect(self, slide, slide_index: int) -> float:
        """
        Title and Content is the default handler — it matches any slide
        that has a title and body content. It should have LOW priority
        so more specific handlers (cover, divider, thank you) win first.
        """
        # Never match first slide (that's the cover)
        if slide_index == 0:
            return 0.0

        texts = extract_text_elements(slide)
        images = extract_images(slide)

        if not texts:
            return 0.0

        # Filter noise
        meaningful = [
            t for t in texts
            if not any(p in t["text"].lower() for p in self.PLACEHOLDER_NOISE)
            and len(t["text"].strip()) > 1
        ]

        if not meaningful:
            return 0.0

        # Basic heuristic: has at least one text element that looks like
        # a title (shorter text) and at least some body content
        has_short_text = any(len(t["text"]) < 80 for t in meaningful)
        has_body_text = any(len(t["text"]) > 30 for t in meaningful)

        # Content slides typically have more text than dividers
        total_text_len = sum(len(t["text"]) for t in meaningful)

        if has_short_text and has_body_text and total_text_len > 50:
            return 0.4  # Low confidence — acts as fallback

        if len(meaningful) >= 2:
            return 0.35

        return 0.1

    # --- Extraction ---

    def extract_content(self, slide, slide_index: int) -> dict:
        """
        Extract title, optional subtitle, and body content from a slide.

        Strategy:
        1. Use source placeholder indices if available — idx 0/3/15 = title,
           idx 1/10 = content body. This is the most reliable signal.
        2. Fall back to font size / position heuristics for non-placeholder shapes.
        3. For the content body placeholder, extract rich paragraph data
           (levels, runs, bold/italic) so formatting can be replayed on output.
        """
        shapes = extract_shapes_with_text(slide)

        result = {
            "title": "",
            "subtitle": "",
            "content": "",
            "rich_paragraphs": [],  # Rich paragraph data for formatted output
            "footer": "",
        }

        if not shapes:
            return result

        # Filter out noise, bare slide numbers, and image license text
        filtered = []
        for s in shapes:
            text = s["text"].strip()
            if not text:
                continue
            if any(p in text.lower() for p in self.PLACEHOLDER_NOISE):
                continue
            if re.match(r"^\d{1,3}$", text) and len(text) <= 3:
                continue
            if self._is_image_caption(text):
                continue
            filtered.append(s)

        if not filtered:
            return result

        # --- Separate by role using placeholder indices ---
        title_shape = None
        subtitle_shape = None
        body_shapes = []
        body_placeholder = None  # The actual placeholder object for rich extraction
        footer_text = ""

        for s in filtered:
            if s["is_placeholder"]:
                idx = s["shape_idx"]
                if idx in (0, 3, 15):
                    # TITLE / CENTER_TITLE
                    title_shape = s
                elif idx in (1, 10):
                    # BODY / CONTENT
                    body_shapes.append(s)
                    # Find the actual placeholder object for rich extraction
                    if body_placeholder is None:
                        for shape in slide.shapes:
                            if (hasattr(shape, 'is_placeholder') and shape.is_placeholder
                                    and shape.placeholder_format.idx == idx):
                                body_placeholder = shape
                                break
                elif idx in (17,):
                    # FOOTER
                    if not footer_text:
                        footer_text = s["text"]
                elif idx in (18,):
                    # SLIDE NUMBER — skip
                    pass
                else:
                    # Unknown placeholder — treat as body
                    body_shapes.append(s)
            elif self._is_footer_text(s["text"]):
                if not footer_text:
                    footer_text = s["text"]
            else:
                # Non-placeholder shapes go to body
                body_shapes.append(s)

        result["footer"] = footer_text

        # --- Fallback title detection if no placeholder found ---
        if not title_shape and body_shapes:
            # Priority 1: Shape with largest font size
            with_font = [s for s in body_shapes if s["font_size"]]
            if with_font:
                candidate = max(with_font, key=lambda s: s["font_size"])
                # Only promote to title if it's reasonably short
                if len(candidate["text"]) < 120:
                    title_shape = candidate
                    body_shapes = [s for s in body_shapes if s is not title_shape]

            # Priority 2: Topmost short shape
            if not title_shape:
                sorted_by_top = sorted(body_shapes, key=lambda s: s["top"])
                for s in sorted_by_top:
                    if len(s["text"]) < 120:
                        title_shape = s
                        body_shapes = [bs for bs in body_shapes if bs is not s]
                        break

        if title_shape:
            result["title"] = title_shape["text"]

        # --- Identify subtitle ---
        # Sort body shapes by vertical position
        body_shapes.sort(key=lambda s: (s["top"], s["left"]))

        if body_shapes and title_shape:
            candidate = body_shapes[0]
            is_short = len(candidate["text"]) < 100
            is_near_top = candidate["top"] < (title_shape["top"] + title_shape["height"] * 2)
            has_more_content = len(body_shapes) > 1

            if is_short and is_near_top and has_more_content:
                result["subtitle"] = candidate["text"]
                body_shapes = body_shapes[1:]

        # --- Body content ---
        if body_shapes:
            result["content"] = "\n".join(s["text"] for s in body_shapes)

        # --- Rich paragraph extraction from the body placeholder ---
        if body_placeholder:
            raw_paras = extract_rich_paragraphs(body_placeholder)

            # If a subtitle was detected, it came from the first paragraph(s)
            # of the body placeholder. Remove those from rich_paragraphs to
            # avoid duplicating text (subtitle goes in idx=31, body in idx=10).
            if result["subtitle"] and raw_paras:
                subtitle_text = result["subtitle"].strip()
                # Skip leading paragraphs that match the subtitle text
                while raw_paras and (
                    raw_paras[0]["is_empty"]
                    or raw_paras[0]["text"].strip() == subtitle_text
                ):
                    removed = raw_paras.pop(0)
                    if removed["text"].strip() == subtitle_text:
                        break  # Found and removed the subtitle para

            result["rich_paragraphs"] = raw_paras

        # --- Append non-placeholder body text (e.g. from group shapes) ---
        # Group shapes and free text boxes contribute to 'content' but aren't
        # in rich_paragraphs (which comes from the body placeholder only).
        # Append them as additional bullet-level paragraphs.
        if body_placeholder and body_shapes:
            placeholder_text = body_placeholder.text_frame.text.strip()
            for s in body_shapes:
                # Skip the body placeholder itself (already in rich_paragraphs)
                if s["is_placeholder"] and s["shape_idx"] in (1, 10):
                    continue
                shape_text = s["text"].strip()
                if shape_text and shape_text not in placeholder_text:
                    for line in shape_text.split("\n"):
                        line = line.strip()
                        if line:
                            result["rich_paragraphs"].append({
                                "level": 0,
                                "runs": [{"text": line, "bold": False, "italic": None}],
                                "text": line,
                                "has_bullet": True,
                                "bullet_char": "•",
                                "is_empty": False,
                            })

        return result

    def _is_image_caption(self, text: str) -> bool:
        """Check if text is an image license/caption that should be filtered."""
        caption_patterns = [
            r"(?i)image\s+licensed\s+through",
            r"(?i)adobe\s+stock[:\s]+\d",
            r"(?i)shutterstock[:\s]+\d",
            r"(?i)getty\s+images",
            r"(?i)©\s*\d{4}",
            r"(?i)source:\s*(http|www)",
        ]
        return any(re.search(p, text) for p in caption_patterns)

    def _is_footer_text(self, text: str) -> bool:
        """Check if text looks like a footer/programme name."""
        footer_patterns = [
            r"(?i)cricos", r"(?i)hbis\s+innovation",
            r"(?i)executive\s+education",
            r"(?i)presentation\s+title",
        ]
        return any(re.search(p, text) for p in footer_patterns)

    # --- Output ---

    # --- Template level mapping ---
    # The UQ template master defines:
    #   Level 0 (lvl1pPr): buNone, bold, accent1 colour — heading/label
    #   Level 1 (lvl2pPr): buNone, regular, tx1 — sub-heading
    #   Level 2 (lvl3pPr): buChar '•', regular, tx1 — bullet point
    #   Level 3 (lvl4pPr): buChar '–', regular, tx1 — sub-bullet
    #   Level 4 (lvl5pPr): buChar '–', regular, tx1 — sub-sub-bullet
    TEMPLATE_HEADING_LEVEL = 0
    TEMPLATE_BULLET_LEVEL = 2
    TEMPLATE_SUBBULLET_LEVEL = 3

    def fill_slide(self, slide, content: dict) -> None:
        """
        Fill Title and Content placeholders.
        If no subtitle, shifts content body up to reclaim the space.
        Uses rich paragraph data when available to preserve formatting.
        """
        from pptx.util import Emu

        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Title
        if content.get("title") and self.PH_TITLE in placeholders:
            placeholders[self.PH_TITLE].text = content["title"]

        # Subtitle
        has_subtitle = bool(content.get("subtitle"))
        if has_subtitle and self.PH_SUBTITLE in placeholders:
            placeholders[self.PH_SUBTITLE].text = content["subtitle"]
        elif not has_subtitle and self.PH_SUBTITLE in placeholders:
            # Clear subtitle and shift content up
            placeholders[self.PH_SUBTITLE].text = ""

        # Content body
        if content.get("content") and self.PH_CONTENT in placeholders:
            ph = placeholders[self.PH_CONTENT]

            # Shift content up if no subtitle.
            if not has_subtitle:
                orig_left = ph.left
                orig_width = ph.width
                new_top = ph.top - self.SUBTITLE_HEIGHT_EMU
                new_height = ph.height + self.SUBTITLE_HEIGHT_EMU
                ph.top = new_top
                ph.left = orig_left
                ph.width = orig_width
                ph.height = new_height

            rich_paras = content.get("rich_paragraphs", [])
            if rich_paras:
                self._fill_rich_content(ph, rich_paras)
            else:
                # Fallback: plain text (for backwards compat)
                ph.text = content["content"]
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = False

        # Footer
        if content.get("footer") and self.PH_FOOTER in placeholders:
            placeholders[self.PH_FOOTER].text = content["footer"]

    def _fill_rich_content(self, placeholder, rich_paras: list[dict]) -> None:
        """
        Fill a content placeholder using rich paragraph data, preserving
        bullet levels, bold/italic formatting, and paragraph structure.

        Mapping strategy (source → template):
        - Source level 0 with no bullet + short text → Template level 0 (heading)
        - Source level 0 with bullet or long text → Template level 2 (bullet)
        - Source level 1+ with bullet → Template level 2 or 3 (bullet/sub-bullet)
        - Source level 2+ → Template level 3 (sub-bullet)
        - Empty paragraphs → preserved as spacers (level 0, no bullet)
        """
        from pptx.oxml.ns import qn

        tf = placeholder.text_frame

        # Clear existing content — remove all but the first paragraph,
        # then we'll overwrite the first and append the rest.
        while len(tf.paragraphs) > 1:
            p_elem = tf.paragraphs[-1]._p
            p_elem.getparent().remove(p_elem)

        # Filter out trailing empty paragraphs (source artefacts)
        while rich_paras and rich_paras[-1]["is_empty"]:
            rich_paras.pop()

        for i, rp in enumerate(rich_paras):
            if i == 0:
                para = tf.paragraphs[0]
                # Clear existing runs from the template paragraph
                for r in list(para._p.findall(qn('a:r'))):
                    para._p.remove(r)
                # Also remove any existing a:br elements
                for br in list(para._p.findall(qn('a:br'))):
                    para._p.remove(br)
            else:
                para = self._add_paragraph(tf)

            # Determine the target template level
            target_level = self._map_source_level(rp)

            # Set paragraph level
            para.level = target_level

            if rp["is_empty"]:
                # Spacer paragraph — keep empty at level 0
                para.level = 0
                # Ensure no bullet on empty lines
                self._set_no_bullet(para)
                continue

            # Populate runs
            if rp["runs"]:
                for run_data in rp["runs"]:
                    run = para.add_run()
                    run.text = run_data["text"]

                    # Apply formatting — only set explicit overrides.
                    # Level 0 in the template is bold by default (from master),
                    # so we DON'T override bold for heading-level paragraphs.
                    if target_level == self.TEMPLATE_HEADING_LEVEL:
                        # Heading: let master bold apply. If source had
                        # explicit bold=False on a run, respect that.
                        if run_data["bold"] is False:
                            run.font.bold = False
                        # Otherwise leave as None (inherit master bold)
                    else:
                        # Bullet/sub-bullet levels: master is not bold.
                        # Preserve source bold where explicitly set.
                        if run_data["bold"] is True:
                            run.font.bold = True
                        # bold=None or False → leave as inherited (not bold)

                    if run_data["italic"] is True:
                        run.font.italic = True
            else:
                # No runs but has text (e.g. from soft-breaks)
                run = para.add_run()
                run.text = rp["text"]
                if target_level != self.TEMPLATE_HEADING_LEVEL:
                    pass  # Inherit template defaults

    def _map_source_level(self, rp: dict) -> int:
        """
        Map a source paragraph to the appropriate UQ template level.

        The UQ template defines:
            Level 0: Heading (bold, purple, no bullet)
            Level 2: Bullet point (•)
            Level 3: Sub-bullet (–)

        Source signals:
            - has_bullet=True → it's a bulleted item
            - has_bullet=False (buNone) → no bullet (heading or plain text)
            - has_bullet=None → inherited from layout (assume bulleted if level > 0)
            - level > 0 → indented (sub-item)
        """
        src_level = rp["level"]
        has_bullet = rp["has_bullet"]
        text = rp["text"].strip()

        if rp["is_empty"]:
            return 0

        # Explicit bullet in source → map to template bullet levels
        if has_bullet is True:
            if src_level == 0:
                return self.TEMPLATE_BULLET_LEVEL
            elif src_level == 1:
                return self.TEMPLATE_BULLET_LEVEL
            else:
                return self.TEMPLATE_SUBBULLET_LEVEL

        # Explicit no-bullet (buNone) in source
        if has_bullet is False:
            # Short text at level 0 → heading
            if src_level == 0 and len(text) < 120:
                return self.TEMPLATE_HEADING_LEVEL
            # Longer text → still heading but could be content
            return self.TEMPLATE_HEADING_LEVEL

        # Inherited bullet (has_bullet is None)
        if src_level == 0:
            # Level 0 with inherited formatting — look at the first run's
            # bold state to decide heading vs bullet.
            # bold=None means inheriting from master, which in academic
            # slides typically means the master has bold on level-0.
            # bold=True is explicitly bold. Either indicates a heading.
            # bold=False is explicitly not bold → likely a bullet item.
            runs = rp.get("runs", [])
            first_run_bold = runs[0].get("bold") if runs else None
            if first_run_bold is None or first_run_bold is True:
                # Inherits or explicitly bold → heading
                return self.TEMPLATE_HEADING_LEVEL
            return self.TEMPLATE_BULLET_LEVEL
        elif src_level <= 2:
            return self.TEMPLATE_BULLET_LEVEL
        else:
            return self.TEMPLATE_SUBBULLET_LEVEL

    @staticmethod
    def _add_paragraph(text_frame):
        """Add a new paragraph to a text frame and return it."""
        from pptx.oxml.ns import qn
        from lxml import etree

        new_p = etree.SubElement(text_frame._txBody, qn('a:p'))
        # Return the last paragraph (the one we just added)
        return text_frame.paragraphs[-1]

    @staticmethod
    def _set_no_bullet(para):
        """Explicitly set buNone on a paragraph to suppress inherited bullets."""
        from pptx.oxml.ns import qn
        from lxml import etree

        pPr = para._p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(para._p, qn('a:pPr'))
            # Insert pPr as the first child
            para._p.insert(0, pPr)
        etree.SubElement(pPr, qn('a:buNone'))

    def get_placeholder_map(self) -> dict:
        return {
            "title": self.PH_TITLE,
            "subtitle": self.PH_SUBTITLE,
            "content": self.PH_CONTENT,
            "footer": self.PH_FOOTER,
            "slide_num": self.PH_SLIDE_NUM,
        }
