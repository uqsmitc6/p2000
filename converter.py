"""
Core conversion engine.

Takes an input PPTX, detects slide types, extracts content,
and produces a brand-compliant output PPTX using the UQ template.

Classification pipeline:
    1. Render all slides to PNG images (once, upfront) if API key provided
    2. Heuristics run first (free, instant)
    3. If confidence >= 0.7 → auto-convert (no API call)
    4. If confidence < 0.7 and API key provided → send slide image to
       Claude Vision for classification
    5. If no API key → fall back to heuristic result with flag

Tiers in the output report:
    - CONVERTED: confident match (heuristic >= 0.7 or API-confirmed)
    - FLAGGED: low confidence, converted but review recommended
    - SKIPPED: no match or classified as "Skip" by the API
"""

import gc
import io
import logging
from pptx import Presentation

from handlers import get_all_handlers, HANDLER_REGISTRY
from utils.template import open_template, add_slide_from_layout, delete_all_original_slides, move_slide_to_position
from utils.references import collect_references, format_references_text, has_meaningful_references
from utils.toc import collect_sections, should_generate_toc, build_toc_content, _has_existing_toc

logger = logging.getLogger("uqslide.converter")

# Confidence thresholds
CONFIDENT_THRESHOLD = 0.7   # No API call needed
CONVERT_THRESHOLD = 0.35    # Minimum to convert without API
# Below CONVERT_THRESHOLD without API → skip


def convert_presentation(
    input_bytes: bytes,
    api_key: str = None,
    model: str = "claude-sonnet-4-6",
    progress_callback=None,
) -> tuple[bytes, dict]:
    """
    Convert an uploaded PPTX to brand-compliant format.

    Args:
        input_bytes: Raw bytes of the uploaded .pptx file
        api_key: Optional Anthropic API key. If provided, Claude Vision
                 classifies ambiguous slides using rendered images.
        model: Claude model for classification (default: claude-sonnet-4-6)
        progress_callback: Optional callable(message: str) for status updates

    Returns:
        (output_bytes, report)
        report includes source_images and output_images dicts for the viewer.
    """
    logger.info("Starting conversion (%d bytes)", len(input_bytes))

    input_prs = Presentation(io.BytesIO(input_bytes))
    handlers = get_all_handlers()
    total_slides = len(input_prs.slides)
    logger.info("Loaded presentation: %d slides", total_slides)

    # Detect programme name from source footers
    programme_name = _detect_programme_name(input_prs)
    if programme_name:
        logger.info("Detected programme name: '%s'", programme_name)

    report = {
        "slides_converted": 0,
        "slides_flagged": 0,
        "slides_skipped": 0,
        "api_calls": 0,
        "errors": [],
        "details": [],
    }

    # --- Pre-render slides to images if API key provided ---
    # Memory guard: skip rendering for very large files (> 8 MB) to avoid
    # exceeding Render free-tier memory limits (512 MB). Heuristic-only
    # conversion still works — just no AI classification or verification.
    MAX_RENDER_BYTES = 8 * 1024 * 1024  # 8 MB
    slide_images = {}
    rendering_skipped_for_memory = False

    if api_key:
        if len(input_bytes) > MAX_RENDER_BYTES:
            logger.warning(
                "Input file too large for rendering (%d MB) — skipping AI "
                "classification/verification to conserve memory",
                len(input_bytes) // (1024 * 1024),
            )
            report["errors"].append(
                f"File too large ({len(input_bytes) // (1024*1024)} MB) for AI-assisted "
                f"classification and verification. Using heuristic-only mode to conserve memory. "
                f"Conversion still works — check results manually."
            )
            rendering_skipped_for_memory = True
            if progress_callback:
                progress_callback(
                    f"Large file ({len(input_bytes) // (1024*1024)} MB) — using heuristic-only mode."
                )
        else:
            if progress_callback:
                progress_callback("Rendering slides to images for AI classification...")
            slide_images, render_diag = _render_slide_images(input_bytes)
            if slide_images:
                logger.info("Rendered %d slide images for Vision classification", len(slide_images))
                if progress_callback:
                    progress_callback(
                        f"Rendered {len(slide_images)} slide images. Classifying..."
                    )
            else:
                logger.warning("Slide rendering failed — falling back to text-only classification")
                report["errors"].append(
                    f"Slide rendering failed: {render_diag}. Used text-only classification fallback."
                )
                if progress_callback:
                    progress_callback(
                        "Could not render slide images. Using text-only fallback."
                    )

    output_prs = open_template()
    new_slides_added = 0

    # --- Pre-scan for existing Acknowledgement of Country slides ---
    # These will be skipped (replaced by auto-inserted branded AoC after verification)
    aoc_handler = handlers.get("Acknowledgement of Country")
    aoc_source_indices = set()
    if aoc_handler:
        for i, s in enumerate(input_prs.slides):
            if aoc_handler.detect(s, i) >= 0.9:
                aoc_source_indices.add(i)

    # --- Collect references & image attributions from source slides ---
    # Must happen before input_prs is deleted. Results stored for compiled
    # references slide generation after verification.
    collected_refs = collect_references(input_prs)
    if has_meaningful_references(collected_refs):
        logger.info(
            "Collected references: %d academic, %d images, %d other",
            len(collected_refs["academic"]),
            len(collected_refs["images"]),
            len(collected_refs["other_sources"]),
        )
    report["collected_references"] = collected_refs

    # --- Pre-scan for Thank You: only allow ONE closing slide ---
    # Find the last slide that looks like a Thank You (heuristic or position)
    # so we can reserve the Thank You layout for only that one.
    last_thank_you_idx = None
    ty_handler = handlers.get("Thank You")
    if ty_handler:
        for i in range(len(input_prs.slides) - 1, max(len(input_prs.slides) - 5, -1), -1):
            if i in aoc_source_indices:
                continue
            s = input_prs.slides[i]
            ty_score = ty_handler.detect(s, i)
            if ty_score >= 0.5:
                last_thank_you_idx = i
                break
        # If heuristics didn't find one, the last slide is the fallback
        if last_thank_you_idx is None:
            last_thank_you_idx = len(input_prs.slides) - 1
    logger.info("Thank You reserved for source slide %d", (last_thank_you_idx or -1) + 1)

    for slide_idx, slide in enumerate(input_prs.slides):

        # Skip source AoC slides — we auto-insert a branded one
        if slide_idx in aoc_source_indices:
            logger.info("Slide %d: skipping existing AoC (will auto-insert branded version)", slide_idx + 1)
            report["details"].append({
                "slide": slide_idx + 1,
                "status": "replaced",  # Not "converted" — no output slide for this source
                "handler": "Acknowledgement of Country",
                "confidence": 0.95,
                "content": None,
                "preview": "Acknowledgement of Country (replaced with branded version)",
                "all_scores": {},
                "classification_method": "auto",
            })
            continue

        if progress_callback:
            progress_callback(
                f"Processing slide {slide_idx + 1} of {total_slides}..."
            )

        # --- Step 1: Heuristic scoring ---
        scores = {}
        for name, handler in handlers.items():
            try:
                scores[name] = handler.detect(slide, slide_idx)
            except Exception as e:
                logger.error("Slide %d: handler '%s' detect() failed: %s", slide_idx + 1, name, e)
                scores[name] = 0.0

        best_name = max(scores, key=scores.get)
        best_confidence = scores[best_name]
        classification_method = "heuristic"
        logger.debug("Slide %d: heuristic → %s (%.2f) | scores=%s",
                      slide_idx + 1, best_name, best_confidence, scores)

        # --- Step 2: API fallback for low-confidence slides ---
        if best_confidence < CONFIDENT_THRESHOLD and api_key:
            logger.info("Slide %d: low confidence (%.2f), calling Vision API", slide_idx + 1, best_confidence)
            api_result = _classify_with_api(
                slide, slide_idx, total_slides, api_key, model,
                slide_image=slide_images.get(slide_idx),
            )

            if api_result and api_result.get("type"):
                report["api_calls"] += 1
                api_type = api_result["type"]
                api_confidence = api_result.get("confidence", 0.8)
                classification_method = "api"
                logger.info("Slide %d: API → %s (%.2f) reason=%s",
                            slide_idx + 1, api_type, api_confidence,
                            api_result.get("reason", ""))

                if api_type == "Skip":
                    # API says skip — but if the heuristic found a viable
                    # match above CONVERT_THRESHOLD, trust the heuristic.
                    # This is especially important when slide rendering failed
                    # (text-only fallback) and the slide has images but no text.
                    if best_confidence >= CONVERT_THRESHOLD:
                        logger.info("Slide %d: API said Skip but heuristic has %s at %.2f — keeping heuristic",
                                    slide_idx + 1, best_name, best_confidence)
                        # Don't skip — fall through to Step 3 with heuristic result
                    else:
                        report["slides_skipped"] += 1
                        report["details"].append({
                            "slide": slide_idx + 1,
                            "status": "skipped",
                            "handler": None,
                            "confidence": api_confidence,
                            "content": None,
                            "preview": _get_slide_preview(slide),
                            "all_scores": scores,
                            "classification_method": "api",
                            "api_reason": api_result.get("reason", ""),
                        })
                        continue

                elif api_type in HANDLER_REGISTRY:
                    # Guard: only allow Thank You for the reserved last slot
                    if api_type == "Thank You" and slide_idx != last_thank_you_idx:
                        logger.info("Slide %d: API suggested Thank You but reserved for slide %d — using Title and Content",
                                    slide_idx + 1, (last_thank_you_idx or -1) + 1)
                        best_name = "Title and Content"
                        best_confidence = api_confidence
                    # Guard: Title Only should not be used when there's
                    # significant body text — the content would be lost.
                    # Downgrade to Title and Content to preserve text.
                    elif api_type == "Title Only" and _slide_has_body_text(slide, min_chars=40):
                        logger.info("Slide %d: API suggested Title Only but slide has body text — using Title and Content",
                                    slide_idx + 1)
                        best_name = "Title and Content"
                        best_confidence = api_confidence
                    else:
                        # API classified it — use that handler
                        best_name = api_type
                        best_confidence = api_confidence
            elif api_result and api_result.get("error"):
                err_msg = f"Slide {slide_idx + 1}: API error — {api_result['error']}"
                logger.error(err_msg)
                report["errors"].append(err_msg)

        # --- Guard: Thank You only for the reserved last slot ---
        if best_name == "Thank You" and slide_idx != last_thank_you_idx:
            logger.info("Slide %d: downgrading Thank You to Title and Content (reserved for slide %d)",
                        slide_idx + 1, (last_thank_you_idx or -1) + 1)
            best_name = "Title and Content"
            best_confidence = scores.get("Title and Content", 0.40)

        # --- Step 3: Convert or skip ---
        best_handler = handlers.get(best_name)
        slide_preview = _get_slide_preview(slide)

        if best_handler and best_confidence >= CONVERT_THRESHOLD:
            try:
                content = best_handler.extract_content(slide, slide_idx)

                # Some handlers dynamically choose their layout based on content
                if hasattr(best_handler, 'get_layout_index'):
                    layout_idx = best_handler.get_layout_index(content)
                else:
                    layout_idx = best_handler.layout_index

                new_slide = add_slide_from_layout(output_prs, layout_idx)
                best_handler.fill_slide(new_slide, content)

                # Preserve visual shapes (diagrams, images) from source
                _preserve_visual_shapes(slide, new_slide, best_handler.name)

                # Remove unfilled placeholders (prevents template prompt text showing)
                _cleanup_empty_placeholders(new_slide, best_handler, content)

                # Populate footer and slide number on every converted slide
                _fill_footer_and_slide_num(
                    new_slide, best_handler, programme_name, slide_idx + 1
                )

                new_slides_added += 1

                if best_confidence >= CONFIDENT_THRESHOLD:
                    status = "converted"
                else:
                    status = "flagged"
                    report["slides_flagged"] += 1

                report["slides_converted"] += 1
                logger.info("Slide %d: %s → %s (%.2f, %s)",
                            slide_idx + 1, status, best_handler.name,
                            best_confidence, classification_method)
                report["details"].append({
                    "slide": slide_idx + 1,
                    "status": status,
                    "handler": best_handler.name,
                    "confidence": best_confidence,
                    "content": content,
                    "preview": slide_preview,
                    "all_scores": scores,
                    "classification_method": classification_method,
                })
            except Exception as e:
                err_msg = f"Slide {slide_idx + 1}: conversion failed ({best_handler.name}) — {e}"
                logger.error(err_msg, exc_info=True)
                report["errors"].append(err_msg)
                report["slides_skipped"] += 1
                report["details"].append({
                    "slide": slide_idx + 1,
                    "status": "skipped",
                    "handler": None,
                    "confidence": best_confidence,
                    "content": None,
                    "preview": slide_preview,
                    "all_scores": scores,
                    "classification_method": classification_method,
                    "reason": f"Error: {e}",
                })
        else:
            logger.info("Slide %d: skipped (best='%s' at %.2f)",
                        slide_idx + 1, best_name, best_confidence)
            report["slides_skipped"] += 1
            report["details"].append({
                "slide": slide_idx + 1,
                "status": "skipped",
                "handler": None,
                "confidence": best_confidence,
                "content": None,
                "preview": slide_preview,
                "all_scores": scores,
                "classification_method": classification_method,
                "reason": f"Best match '{best_name}' at {best_confidence:.2f}",
            })

    # Remove original template slides
    if new_slides_added > 0:
        delete_all_original_slides(output_prs, num_new_slides=new_slides_added)

    # --- Assign output_index to each converted/flagged detail ---
    # At this point the output has NO AoC — clean 1:1 source→output mapping.
    output_idx_counter = 0
    for detail in report["details"]:
        if detail["status"] in ("converted", "flagged"):
            detail["output_index"] = output_idx_counter
            output_idx_counter += 1

    # Save pre-AoC output for verification (1:1 mapping, no offset issues)
    pre_aoc_buffer = io.BytesIO()
    output_prs.save(pre_aoc_buffer)
    pre_aoc_buffer.seek(0)
    pre_aoc_bytes = pre_aoc_buffer.getvalue()

    # Free the output Presentation and input Presentation to reclaim RAM
    # before the memory-heavy rendering/verification step.
    del output_prs, input_prs, pre_aoc_buffer
    gc.collect()

    logger.info("Conversion complete (pre-AoC): %d converted, %d flagged, %d skipped, %d API calls, %d errors",
                report["slides_converted"], report["slides_flagged"],
                report["slides_skipped"], report["api_calls"], len(report["errors"]))

    # --- Post-conversion verification ---
    # Send original + branded slide pairs to Claude for QA checking.
    # Runs on pre-AoC output so source→output mapping is clean 1:1.
    output_images = {}
    if api_key and slide_images:
        if progress_callback:
            progress_callback("Verifying converted slides...")

        verification_results, output_images = _verify_all_slides(
            pre_aoc_bytes, slide_images, report, api_key, model,
            progress_callback,
        )
        report["verification"] = verification_results

        # Count issues
        v_passed = sum(1 for v in verification_results if v.get("pass") is True)
        v_issues = sum(1 for v in verification_results if v.get("pass") is False)
        v_errors = sum(1 for v in verification_results if v.get("pass") is None)
        report["verification_summary"] = {
            "total": len(verification_results),
            "passed": v_passed,
            "issues_found": v_issues,
            "errors": v_errors,
        }
        report["api_calls"] += len(verification_results)

        logger.info("Verification: %d passed, %d issues found, %d errors",
                     v_passed, v_issues, v_errors)

        if progress_callback:
            progress_callback(
                f"Verification complete: {v_passed} passed, {v_issues} issues found"
            )

    # --- Store rendered images in report for in-browser viewer ---
    report["source_images"] = slide_images  # {source_slide_index: PNG bytes}
    report["output_images"] = output_images  # {output_slide_index: PNG bytes}

    # --- Insert auto-generated slides (AFTER verification) ---
    # Load a fresh Presentation from pre-AoC bytes to avoid python-pptx
    # internal part counter issues from double-saving the same object.
    final_prs = Presentation(io.BytesIO(pre_aoc_bytes))
    auto_slides_inserted = False

    # 1. Insert Acknowledgement of Country as slide 2
    if aoc_handler:
        try:
            if progress_callback:
                progress_callback("Inserting Acknowledgement of Country...")
            aoc_content = aoc_handler._get_standard_content()
            aoc_slide = add_slide_from_layout(final_prs, aoc_handler.layout_index)
            aoc_handler.fill_slide(aoc_slide, aoc_content)
            _fill_footer_and_slide_num(aoc_slide, aoc_handler, programme_name, 2)
            # Move from end (where add_slide puts it) to position 1 (after cover)
            move_slide_to_position(final_prs, len(final_prs.slides) - 1, 1)
            report["slides_converted"] += 1
            auto_slides_inserted = True
            # Insert AoC detail after cover detail in report
            aoc_detail = {
                "slide": "2 (auto)",
                "status": "converted",
                "handler": "Acknowledgement of Country",
                "confidence": 1.0,
                "content": aoc_content,
                "preview": "Acknowledgement of Country (auto-inserted)",
                "all_scores": {},
                "classification_method": "auto",
            }
            cover_idx = next(
                (i for i, d in enumerate(report["details"])
                 if d.get("handler") == "Cover 1"),
                0,
            )
            report["details"].insert(cover_idx + 1, aoc_detail)
            logger.info("Auto-inserted Acknowledgement of Country as slide 2")
        except Exception as e:
            logger.error("Failed to insert AoC slide: %s", e)
            report["errors"].append(f"AoC auto-insert failed: {e}")

    # 2. Insert Table of Contents (after Cover + AoC, before content)
    toc_sections = collect_sections(final_prs)
    if should_generate_toc(toc_sections) and not _has_existing_toc(final_prs):
        try:
            if progress_callback:
                progress_callback("Generating Table of Contents slide...")
            _insert_toc_slide(final_prs, toc_sections, programme_name)
            report["slides_converted"] += 1
            auto_slides_inserted = True
            toc_detail = {
                "slide": "auto (toc)",
                "status": "converted",
                "handler": "Contents (compiled)",
                "confidence": 1.0,
                "content": {"sections": len(toc_sections)},
                "preview": f"Table of Contents ({len(toc_sections)} sections, auto-generated)",
                "all_scores": {},
                "classification_method": "auto",
            }
            # Insert after AoC detail (or after cover if no AoC)
            aoc_idx = next(
                (i for i, d in enumerate(report["details"])
                 if d.get("handler") == "Acknowledgement of Country"),
                None,
            )
            if aoc_idx is not None:
                report["details"].insert(aoc_idx + 1, toc_detail)
            else:
                cover_idx = next(
                    (i for i, d in enumerate(report["details"])
                     if d.get("handler") == "Cover 1"),
                    0,
                )
                report["details"].insert(cover_idx + 1, toc_detail)
            logger.info("Auto-inserted Table of Contents (%d sections)", len(toc_sections))
        except Exception as e:
            logger.error("Failed to insert ToC slide: %s", e)
            report["errors"].append(f"ToC auto-insert failed: {e}")

    # 3. Insert compiled References & Image Credits slide (before Thank You)
    if has_meaningful_references(collected_refs):
        try:
            if progress_callback:
                progress_callback("Generating References & Image Credits slide...")
            _insert_compiled_references(final_prs, collected_refs, programme_name)
            report["slides_converted"] += 1
            auto_slides_inserted = True
            # Add detail to report
            refs_detail = {
                "slide": "auto (refs)",
                "status": "converted",
                "handler": "References (compiled)",
                "confidence": 1.0,
                "content": {
                    "academic_count": len(collected_refs["academic"]),
                    "image_count": len(collected_refs["images"]),
                    "other_count": len(collected_refs["other_sources"]),
                },
                "preview": "Compiled References & Image Credits (auto-generated)",
                "all_scores": {},
                "classification_method": "auto",
            }
            report["details"].append(refs_detail)
            logger.info(
                "Auto-inserted compiled References slide (%d academic, %d images, %d other)",
                len(collected_refs["academic"]),
                len(collected_refs["images"]),
                len(collected_refs["other_sources"]),
            )
        except Exception as e:
            logger.error("Failed to insert compiled References slide: %s", e)
            report["errors"].append(f"References slide auto-insert failed: {e}")

    # Save final output
    if auto_slides_inserted:
        output_buffer = io.BytesIO()
        final_prs.save(output_buffer)
        output_buffer.seek(0)
        output_bytes = output_buffer.getvalue()
    else:
        output_bytes = pre_aoc_bytes

    del final_prs

    return output_bytes, report


def _verify_all_slides(
    output_bytes: bytes,
    source_images: dict,
    report: dict,
    api_key: str,
    model: str,
    progress_callback=None,
) -> tuple[list, dict]:
    """
    Verify all converted slides by comparing source and output images.

    The output_bytes should be the pre-AoC output (no Acknowledgement of
    Country inserted yet) so source→output mapping is clean 1:1.

    Args:
        output_bytes: branded PPTX bytes (pre-AoC)
        source_images: dict mapping source slide_index → PNG bytes
        report: conversion report (to read slide details with output_index)
        api_key: Anthropic API key
        model: Claude model to use
        progress_callback: optional status callback

    Returns:
        Tuple of (list of verification results, output_images dict).
        output_images maps output_slide_index → PNG bytes.
    """
    from utils.classifier import verify_slide_pair

    # Render the output slides
    output_images, out_diag = _render_slide_images(output_bytes)
    if not output_images:
        logger.warning("Could not render output slides for verification: %s", out_diag)
        return [], {}

    total_output = len(output_images)
    results = []

    # Walk details — each converted/flagged detail has an output_index
    # assigned earlier. The mapping is now 1:1 (no AoC offset).
    for detail in report["details"]:
        if detail["status"] not in ("converted", "flagged"):
            continue

        source_slide_num = detail["slide"]
        handler_name = detail.get("handler", "Unknown")
        output_idx = detail.get("output_index")

        if not isinstance(source_slide_num, int) or output_idx is None:
            continue

        source_idx = source_slide_num - 1
        source_img = source_images.get(source_idx)
        output_img = output_images.get(output_idx)

        if source_img and output_img:
            if progress_callback:
                progress_callback(
                    f"Verifying slide {output_idx + 1}/{total_output} "
                    f"(source slide {source_slide_num})..."
                )

            v_result = verify_slide_pair(
                source_image=source_img,
                output_image=output_img,
                slide_number=source_slide_num,
                total_slides=total_output,
                handler_name=handler_name,
                api_key=api_key,
                model=model,
            )
            v_result["source_slide"] = source_slide_num
            v_result["output_slide"] = output_idx + 1
            v_result["handler"] = handler_name
            results.append(v_result)

            # Attach to report detail for UI display
            detail["verification"] = v_result
        else:
            logger.debug("Skipping verification for source %d / output %d (missing image)",
                         source_slide_num, output_idx + 1)

    return results, output_images


def _render_slide_images(input_bytes: bytes) -> tuple[dict, str]:
    """
    Render all slides to PNG images using LibreOffice.

    Returns:
        Tuple of (dict mapping slide_index → PNG bytes, diagnostic_message)
        Empty dict if rendering fails.
    """
    try:
        from utils.renderer import render_slides_to_images, is_libreoffice_available
        if not is_libreoffice_available():
            msg = "LibreOffice not available — check Docker image has libreoffice-impress installed"
            logger.error(msg)
            return {}, msg
        logger.info("LibreOffice available, starting render of %d bytes...", len(input_bytes))
        images, diag = render_slides_to_images(input_bytes, dpi=96)
        logger.info("Render complete: %d slide images produced. Diag: %s", len(images), diag)
        return {i: img for i, img in enumerate(images)}, diag
    except Exception as e:
        msg = f"Slide rendering failed: {e}"
        logger.error(msg, exc_info=True)
        return {}, msg


def _classify_with_api(slide, slide_index, total_slides, api_key, model,
                       slide_image=None):
    """Call the AI classifier. Returns result dict or None on failure."""
    try:
        from utils.classifier import classify_slide_with_api
        return classify_slide_with_api(
            slide, slide_index, total_slides, api_key, model,
            slide_image=slide_image,
        )
    except Exception as e:
        logger.error("API classifier failed for slide %d: %s", slide_index + 1, e, exc_info=True)
        return None


def _slide_has_body_text(slide, min_chars: int = 30) -> bool:
    """
    Check if a slide has significant body text (non-title text).
    Used to guard against API misclassifying text-heavy slides as Title Only.
    Counts text from content placeholders, text boxes, and group shapes.
    """
    body_text_len = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip title placeholder (idx 0) and footer/slide number
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            ph_idx = shape.placeholder_format.idx
            if ph_idx in (0, 17, 18):  # title, footer, slide number
                continue
        text = shape.text_frame.text.strip()
        body_text_len += len(text)
    # Also check group shapes for text
    for shape in slide.shapes:
        if shape.shape_type == 6:  # Group
            try:
                from utils.extractor import _extract_group_text
                group_texts = _extract_group_text(shape)
                for text in group_texts:
                    body_text_len += len(text) if isinstance(text, str) else len(text.get("text", ""))
            except Exception:
                pass
    return body_text_len >= min_chars


def _get_slide_preview(slide) -> str:
    """Get a short text preview of a slide for reporting."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) > 1:
                texts.append(t[:60])
    return " | ".join(texts[:3])[:120] if texts else "(no text)"


def _detect_programme_name(prs) -> str:
    """
    Detect the programme/course name for footer text.

    Strategy:
    1. Check footer placeholders (idx 17) on slides 2-5 for an existing
       programme name like "Executive Education" or "Negotiating For Success".
    2. Fall back to the cover slide title (largest text on slide 1).
    """
    if len(prs.slides) == 0:
        return ""

    # Strategy 1: Look for footer placeholder text on early slides
    for i in range(1, min(6, len(prs.slides))):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if (hasattr(shape, 'is_placeholder') and shape.is_placeholder
                    and shape.placeholder_format.idx == 17):
                text = shape.text_frame.text.strip()
                if text and len(text) < 80:
                    return text

    # Strategy 2: Cover slide title (largest font)
    cover = prs.slides[0]

    # Check placeholder idx 0 first
    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        if (hasattr(shape, 'is_placeholder') and shape.is_placeholder
                and shape.placeholder_format.idx in (0, 3, 15)):
            text = shape.text_frame.text.strip()
            if text and len(text) < 120:
                return text

    # Check shape named "Title" (some slides have non-placeholder titles)
    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        if "title" in shape.name.lower():
            text = shape.text_frame.text.strip()
            if text and len(text) < 120:
                return text

    # Fallback: largest font on the first slide
    best_text = ""
    best_size = 0
    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size.pt > best_size:
                    best_size = run.font.size.pt
                    best_text = shape.text_frame.text.strip()

    return best_text if len(best_text) < 120 else ""


def _cleanup_empty_placeholders(slide, handler, content: dict):
    """
    Remove unfilled placeholders so template prompt text doesn't show through.

    Strategy: iterate through ALL placeholders on the slide and remove any
    that still contain only template prompt text (unfilled). This is safer
    than matching against the content dict keys, because handler content
    keys don't always match placeholder map keys (e.g. Split Content uses
    "narrow_content"/"wide_content" but the map has "left"/"right").

    Skips title, footer, and slide_num (always handled separately).
    Also skips placeholders that have real content (text, tables, images).
    """
    # Placeholder indices to always keep
    ph_map = handler.get_placeholder_map()
    keep_indices = set()
    for key in ("title", "footer", "slide_num"):
        idx = ph_map.get(key)
        if isinstance(idx, int):
            keep_indices.add(idx)

    # Template prompt text patterns — these indicate unfilled placeholders
    PROMPT_PATTERNS = [
        "click to add", "click to edit", "click icon",
        "[click", "[add", "[your", "[insert", "[subtitle",
    ]

    for ph in list(slide.placeholders):
        ph_idx = ph.placeholder_format.idx

        # Always keep title, footer, slide number
        if ph_idx in keep_indices:
            continue

        # Check if placeholder has real content
        has_content = False

        # Check for table (GraphicFrame with table)
        try:
            if ph.has_table:
                has_content = True
        except (AttributeError, TypeError):
            pass

        # Check for chart
        try:
            if ph.has_chart:
                has_content = True
        except (AttributeError, TypeError):
            pass

        # Check for text content
        if not has_content:
            try:
                text = ph.text_frame.text.strip() if ph.has_text_frame else ""
                if text:
                    # Check if it's just template prompt text
                    text_lower = text.lower()
                    is_prompt = any(p in text_lower for p in PROMPT_PATTERNS)
                    if not is_prompt:
                        has_content = True
            except (AttributeError, TypeError):
                pass

        # Check for image
        if not has_content:
            try:
                if hasattr(ph, 'image') and ph.image:
                    has_content = True
            except (ValueError, AttributeError, TypeError):
                pass

        if not has_content:
            # Remove the empty/prompt-only placeholder
            sp = ph._element
            sp.getparent().remove(sp)


def _preserve_visual_shapes(source_slide, output_slide, handler_name: str):
    """
    Copy non-placeholder visual shapes (group shapes, images) from the
    source slide to the output slide.  This preserves diagrams, charts,
    and images that the handler's text-based fill_slide can't reproduce.

    Also recovers images from source placeholders when the output
    placeholder was overwritten with text (e.g. a picture placeholder
    in the source that becomes a text placeholder in the branded output).

    Skips:
    - Text placeholders (already handled by the handler)
    - Very large background images (> 90% of slide in both dimensions)
    - Very small shapes (logos, icons < 8% of slide width)
    - TextImage handler slides (images are already placed via placeholder)

    Uses XML deep-copy for group shapes and SVG/EMF pictures.
    Uses python-pptx add_picture for standard image shapes.
    """
    import io
    from copy import deepcopy
    from pptx.util import Emu

    # Handlers that place images via picture placeholder — skip to avoid dupes
    if handler_name in ("Text with Image", "Text with Image Alt", "Picture with Caption",
                         "Top Image + Content", "Picture with Pullout", "Image Collage",
                         "Text with 4 Images", "Three Column Text & Images",
                         "Acknowledgement of Country", "Quote", "Quote 2"):
        return

    # Handlers that extract body text into placeholders — only preserve
    # images, NOT group shapes (group text is already in rich_paragraphs
    # and the visual group would overlap with the template's content area).
    # EXCEPTION: if the slide has very little non-group text, the "body
    # text" was likely extracted FROM the group shape and the diagram is
    # the primary content. In that case, preserve group shapes.
    TEXT_HANDLERS = {
        "Title and Content", "Two Content", "Split Content",
        "References", "Quote",
    }
    skip_groups = handler_name in TEXT_HANDLERS

    # Check if the slide is primarily a diagram (group shapes hold the
    # main content). If so, don't skip groups even for TEXT_HANDLERS.
    if skip_groups:
        non_group_text_len = 0
        has_groups = False
        for shape in source_slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                continue
            if shape.shape_type == 6:  # Group
                has_groups = True
            elif shape.has_text_frame:
                non_group_text_len += len(shape.text_frame.text.strip())
        # If there are group shapes but very little non-group text,
        # the diagram IS the content — preserve it
        if has_groups and non_group_text_len < 100:
            skip_groups = False
            logger.info("  Preserving group shapes (diagram-dominant slide)")

    SLIDE_W = 12192000  # EMU
    SLIDE_H = 6858000

    shapes_copied = 0

    for shape in source_slide.shapes:
        # Skip placeholders — handled separately below
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            continue

        shape_type = shape.shape_type

        # Get dimensions for filtering
        w = shape.width or 0
        h = shape.height or 0
        w_pct = w / SLIDE_W if SLIDE_W else 0
        h_pct = h / SLIDE_H if SLIDE_H else 0

        # Skip very large shapes (full-slide backgrounds/decorative fills)
        # Raised from 70%/50% — 76%×71% content images were being filtered
        if w_pct > 0.90 and h_pct > 0.85:
            continue

        # Skip very small shapes (logos, tiny icons)
        if w_pct < 0.08 and h_pct < 0.08:
            continue

        # Group shapes (shape_type 6) — deep copy XML
        # Only for handlers without body text placeholders (e.g. Title Only)
        if shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            if skip_groups:
                continue
            _copy_group_shape(source_slide, output_slide, shape)
            shapes_copied += 1

        # Picture shapes (shape_type 13) — copy via add_picture
        elif shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                blob = shape.image.blob
                image_stream = io.BytesIO(blob)
                output_slide.shapes.add_picture(
                    image_stream, shape.left, shape.top,
                    shape.width, shape.height
                )
                shapes_copied += 1
            except Exception as e:
                # Fallback for SVG/EMF images: python-pptx can't extract
                # the blob when the main a:blip has no r:embed (SVG stores
                # the ref in asvg:svgBlip extension). Use XML deep-copy.
                logger.debug("  add_picture failed for '%s': %s — trying XML copy", shape.name, e)
                try:
                    _copy_picture_shape_xml(source_slide, output_slide, shape)
                    shapes_copied += 1
                except Exception as e2:
                    logger.warning("  Could not preserve picture '%s': %s", shape.name, e2)

        # Shapes with image fill (e.g. rectangles with picture fills)
        elif hasattr(shape, 'image'):
            try:
                blob = shape.image.blob
                image_stream = io.BytesIO(blob)
                output_slide.shapes.add_picture(
                    image_stream, shape.left, shape.top,
                    shape.width, shape.height
                )
                shapes_copied += 1
            except Exception as e:
                logger.warning("  Could not preserve image-fill shape '%s': %s", shape.name, e)

        # Visual shapes: freeforms, lines, auto shapes, connectors,
        # SmartArt diagrams — XML deep copy to preserve diagram layouts.
        # These are non-text visual elements that can't be reproduced
        # from text extraction alone.
        # Shape types: 1=AUTO_SHAPE, 5=FREEFORM, 9=LINE, None=SmartArt/diagram
        elif shape_type in (1, 5, 9) or shape_type is None:
            # For text handlers, only copy these if the slide is
            # diagram-dominant (i.e. skip_groups is already False)
            if skip_groups:
                # But check: if the shape is purely visual (line, freeform
                # with no/little text), always preserve it
                shape_text = ""
                try:
                    if shape.has_text_frame:
                        shape_text = shape.text_frame.text.strip()
                except (AttributeError, TypeError):
                    pass
                # Lines and freeforms with no/minimal text are visual elements
                if shape_type in (5, 9) and len(shape_text) < 10:
                    pass  # Allow — these are visual diagram elements
                elif shape_type is None:
                    pass  # SmartArt — always preserve
                else:
                    continue  # Skip text-heavy auto shapes for text handlers
            try:
                _copy_shape_xml(source_slide, output_slide, shape)
                shapes_copied += 1
            except Exception as e:
                logger.warning("  Could not preserve shape '%s' (type=%s): %s",
                               shape.name, shape_type, e)

    # --- Recover images from source placeholders ---
    # When a source placeholder contains an image (e.g. a picture
    # placeholder) but the handler overwrites it with text in the output,
    # the image is lost. Detect and re-add these as free-standing pictures.
    shapes_copied += _recover_placeholder_images(
        source_slide, output_slide, handler_name
    )

    if shapes_copied:
        logger.info("  Preserved %d visual shape(s) from source slide", shapes_copied)


def _copy_shape_xml(source_slide, output_slide, shape):
    """
    Generic XML deep-copy for any shape (freeforms, lines, auto shapes,
    SmartArt diagrams, etc.). Transfers all embedded relationships.
    """
    from copy import deepcopy

    new_sp = deepcopy(shape._element)

    src_part = source_slide.part
    dst_part = output_slide.part

    R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

    for elem in new_sp.iter():
        for attr_name in (f'{R_NS}embed', f'{R_NS}link', f'{R_NS}id'):
            old_rId = elem.get(attr_name)
            if old_rId:
                try:
                    rel = src_part.rels[old_rId]
                    new_rId = dst_part.relate_to(rel.target_part, rel.reltype)
                    elem.set(attr_name, new_rId)
                except (KeyError, Exception):
                    pass

    output_slide.shapes._spTree.append(new_sp)


def _copy_picture_shape_xml(source_slide, output_slide, pic_shape):
    """
    Deep-copy a picture shape via XML, transferring all relationships.
    Used as fallback for SVG/EMF images where python-pptx .image.blob
    fails (no r:embed on the main a:blip element).
    """
    from copy import deepcopy

    new_sp = deepcopy(pic_shape._element)

    src_part = source_slide.part
    dst_part = output_slide.part

    R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

    for elem in new_sp.iter():
        for attr_name in (f'{R_NS}embed', f'{R_NS}link'):
            old_rId = elem.get(attr_name)
            if old_rId:
                try:
                    rel = src_part.rels[old_rId]
                    new_rId = dst_part.relate_to(rel.target_part, rel.reltype)
                    elem.set(attr_name, new_rId)
                except (KeyError, Exception) as e:
                    logger.debug("  Skipped broken rel %s: %s", old_rId, e)

    output_slide.shapes._spTree.append(new_sp)


def _recover_placeholder_images(source_slide, output_slide, handler_name: str) -> int:
    """
    Check source placeholders for images that are missing from the output.

    When the classifier picks a text-based layout for a slide whose source
    had an image in a content placeholder, the handler fills that placeholder
    with text and the image is lost. This function detects those cases and
    copies the image as a free-standing picture on the output slide.

    Returns the number of images recovered.
    """
    import io

    # Only relevant for text-based handlers
    TEXT_HANDLERS = {
        "Title and Content", "Two Content", "Split Content",
        "References", "Title Only",
    }
    if handler_name not in TEXT_HANDLERS:
        return 0

    recovered = 0

    for src_ph in source_slide.placeholders:
        idx = src_ph.placeholder_format.idx
        # Skip title (0), subtitle (1), footer (17), slide number (18)
        if idx in (0, 1, 17, 18):
            continue

        # Check if source placeholder has an image
        src_has_image = False
        try:
            if hasattr(src_ph, 'image') and src_ph.image:
                src_has_image = True
        except (ValueError, AttributeError, TypeError):
            pass

        if not src_has_image:
            continue

        # Check if the output placeholder for this idx has an image
        out_has_image = False
        for out_ph in output_slide.placeholders:
            if out_ph.placeholder_format.idx == idx:
                try:
                    if hasattr(out_ph, 'image') and out_ph.image:
                        out_has_image = True
                except (ValueError, AttributeError, TypeError):
                    pass
                break

        if out_has_image:
            continue

        # Source had image, output doesn't — copy it
        try:
            blob = src_ph.image.blob
            image_stream = io.BytesIO(blob)
            output_slide.shapes.add_picture(
                image_stream, src_ph.left, src_ph.top,
                src_ph.width, src_ph.height
            )
            recovered += 1
            logger.info("  Recovered placeholder image (ph=%d) from source", idx)
        except Exception as e:
            # SVG/EMF fallback
            try:
                _copy_picture_shape_xml(source_slide, output_slide, src_ph)
                recovered += 1
                logger.info("  Recovered placeholder image via XML (ph=%d)", idx)
            except Exception as e2:
                logger.warning("  Could not recover placeholder image (ph=%d): %s", idx, e2)

    return recovered


def _copy_group_shape(source_slide, output_slide, group_shape):
    """
    Deep copy a group shape from source to output slide.

    Group shapes may contain embedded images (via relationships).
    We copy the XML and re-create any image relationships on the output slide.
    """
    from copy import deepcopy
    from lxml import etree

    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    # Deep copy the group shape XML
    new_sp = deepcopy(group_shape._element)

    # Find all relationship references (r:embed, r:link) within the group
    # and re-create them on the output slide
    src_part = source_slide.part
    dst_part = output_slide.part

    for elem in new_sp.iter():
        for attr_name in ('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed',
                          '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link'):
            old_rId = elem.get(attr_name)
            if old_rId:
                try:
                    rel = src_part.rels[old_rId]
                    # Copy the target part to the output slide
                    new_rId = dst_part.relate_to(rel.target_part, rel.reltype)
                    elem.set(attr_name, new_rId)
                except (KeyError, Exception):
                    pass  # Skip broken relationships

    # Append the copied group to the output slide's shape tree
    output_slide.shapes._spTree.append(new_sp)


def _insert_compiled_references(
    prs,
    refs: dict,
    programme_name: str,
) -> None:
    """
    Insert a compiled References & Image Credits slide into the presentation.

    Uses the "Title and Content" layout (index 6). The slide is inserted
    before the last slide if that slide is a Thank You, otherwise appended
    at the end.

    CRITICAL: Never set font properties — inherit from template.
    """
    import re as _re

    # Build the references text
    parts = []

    # Academic references
    if refs["academic"]:
        for ref in refs["academic"]:
            parts.append(ref["text"])

    # Image credits — consolidate Adobe Stock IDs
    if refs["images"]:
        if parts:
            parts.append("")  # separator
        adobe_items = [r for r in refs["images"] if r["type"] == "adobe_stock"]
        other_items = [r for r in refs["images"] if r["type"] != "adobe_stock"]

        if adobe_items:
            stock_ids = []
            for r in adobe_items:
                match = _re.search(r"(\d{5,})", r["text"])
                if match:
                    stock_ids.append(match.group(1))
                else:
                    stock_ids.append(r["text"])
            parts.append("Images licensed through Adobe Stock: " + ", ".join(stock_ids))

        for r in other_items:
            parts.append(r["text"])

    # Other sources
    if refs["other_sources"]:
        if parts:
            parts.append("")  # separator
        for ref in refs["other_sources"]:
            parts.append(ref["text"])

    if not parts:
        return

    content_text = "\n".join(parts)

    # Determine title based on what we have
    has_academic = bool(refs["academic"])
    has_images = bool(refs["images"])
    if has_academic and has_images:
        title = "References & Image Credits"
    elif has_academic:
        title = "References"
    else:
        title = "Image Credits"

    # Create slide using Title and Content layout (index 6)
    slide = add_slide_from_layout(prs, 6)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Fill title
    if 0 in placeholders:
        placeholders[0].text = title

    # Fill content
    if 10 in placeholders:
        placeholders[10].text = content_text

    # Fill footer
    if 17 in placeholders and programme_name:
        placeholders[17].text = programme_name

    # Check if last slide before this one is Thank You — if so, move refs before it
    total = len(prs.slides)
    if total >= 3:
        # The new refs slide is at position total-1 (last). Check slide before it.
        prev_slide = prs.slides[total - 2]
        prev_text = ""
        for shape in prev_slide.shapes:
            if shape.has_text_frame:
                prev_text += shape.text_frame.text.lower() + " "

        prev_layout = prev_slide.slide_layout.name.lower()
        is_thank_you = (
            "thank you" in prev_layout
            or any(kw in prev_text for kw in [
                "thank you", "thanks", "questions", "contact",
                "execed@", "exceed@",
            ])
        )

        if is_thank_you:
            # Move refs slide from end to just before the Thank You
            move_slide_to_position(prs, total - 1, total - 2)
            logger.info("Moved References slide before Thank You (position %d)", total - 1)


def _insert_toc_slide(
    prs,
    sections: list[dict],
    programme_name: str,
) -> None:
    """
    Insert a Table of Contents slide using the Contents 2 layout (index 4).

    Creates a two-column table with zero-padded numbers and section titles,
    matching the pattern used in human-branded decks (e.g. Climate Finance).

    The slide is inserted after Cover and AoC — typically position 2 or 3.

    CRITICAL: Never set font properties — inherit from template.
    """
    from utils.toc import CONTENTS_2_LAYOUT_INDEX, PH_TITLE, PH_TABLE, PH_FOOTER, PH_SLIDE_NUM, build_toc_content

    rows = build_toc_content(sections)
    if not rows:
        return

    # Create the slide
    slide = add_slide_from_layout(prs, CONTENTS_2_LAYOUT_INDEX)
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Fill title
    if PH_TITLE in placeholders:
        placeholders[PH_TITLE].text = "Contents"

    # Insert table into the table placeholder
    if PH_TABLE in placeholders:
        table_ph = placeholders[PH_TABLE]
        # insert_table returns a GraphicFrame; access .table for the Table object
        graphic_frame = table_ph.insert_table(rows=len(rows), cols=2)
        table = graphic_frame.table

        # Populate cells — ONLY set .text, formatting inherits from master
        for row_idx, (num, title) in enumerate(rows):
            table.cell(row_idx, 0).text = num
            table.cell(row_idx, 1).text = title

    # Fill footer
    if PH_FOOTER in placeholders and programme_name:
        placeholders[PH_FOOTER].text = programme_name

    # Determine insertion position: after Cover (0) and AoC (1), so position 2
    # If AoC was inserted, it's at index 1. ToC goes at index 2.
    # If no AoC, ToC goes at index 1.
    insert_pos = 1  # Default: after cover
    if len(prs.slides) >= 2:
        # Check if slide at index 1 is AoC — check both layout name and title text
        second_slide = prs.slides[1]
        second_layout = second_slide.slide_layout.name.lower()
        second_title = ""
        for shape in second_slide.shapes:
            if (shape.has_text_frame and hasattr(shape, 'is_placeholder')
                    and shape.is_placeholder and shape.placeholder_format.idx == 0):
                second_title = shape.text_frame.text.lower()
                break
        is_aoc = (
            "acknowledgement" in second_layout
            or "country" in second_layout
            or "acknowledgement" in second_title
            or "country" in second_title
        )
        if is_aoc:
            insert_pos = 2

    # The new slide was appended at the end — move it to the correct position
    move_slide_to_position(prs, len(prs.slides) - 1, insert_pos)
    logger.info("Inserted Contents slide at position %d", insert_pos + 1)


def _fill_footer_and_slide_num(slide, handler, programme_name: str, slide_number: int):
    """
    Populate footer and slide number placeholders on a converted slide.
    Uses the handler's placeholder map to find the right indices.
    """
    ph_map = handler.get_placeholder_map()
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Footer (programme name) — always override with detected programme name
    footer_idx = ph_map.get("footer")
    if footer_idx is not None and footer_idx in placeholders and programme_name:
        placeholders[footer_idx].text = programme_name

    # Slide number
    slide_num_idx = ph_map.get("slide_num")
    if slide_num_idx is not None and slide_num_idx in placeholders:
        placeholders[slide_num_idx].text = str(slide_number)


def convert_cover_only(input_bytes: bytes) -> tuple[bytes, dict]:
    """
    Simplified conversion — only processes the first slide as a cover.
    """
    input_prs = Presentation(io.BytesIO(input_bytes))

    if len(input_prs.slides) == 0:
        raise ValueError("The uploaded file contains no slides.")

    from handlers.cover1 import Cover1Handler
    handler = Cover1Handler()

    slide = input_prs.slides[0]
    content = handler.extract_content(slide, 0)

    output_prs = open_template()
    new_slide = add_slide_from_layout(output_prs, handler.layout_index)
    handler.fill_slide(new_slide, content)
    delete_all_original_slides(output_prs, num_new_slides=1)

    output_buffer = io.BytesIO()
    output_prs.save(output_buffer)
    output_buffer.seek(0)

    report = {
        "slides_converted": 1,
        "slides_flagged": 0,
        "slides_skipped": len(input_prs.slides) - 1,
        "api_calls": 0,
        "details": [{
            "slide": 1,
            "status": "converted",
            "handler": handler.name,
            "content": content,
        }],
    }

    return output_buffer.getvalue(), report
