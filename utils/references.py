"""
References & Image Attribution Collector.

Scans all slides in a converted presentation to collect:
1. Academic references (APA-style citations, DOIs, journal names)
2. Image attributions (Adobe Stock IDs, Unsplash/Pexels credits, etc.)
3. Other source attributions ("Source: ...", "Diagram created by ...")

Produces a deduplicated, categorised set of references that can be
compiled into a summary References & Image Credits slide.

The collector works on the ORIGINAL source slides (not converted output)
because inline attributions are often in small text boxes that may not
survive conversion into template placeholders.
"""

import re
from pptx.presentation import Presentation


# --- Pattern definitions ---

# Academic citation patterns
ACADEMIC_PATTERNS = [
    # Full APA citations: Author (Year). Title. Journal, Vol(Issue), Pages.
    re.compile(
        r"[A-Z][a-z]+,?\s+[A-Z]\..*?\(\d{4}\).*?(?:Journal|Review|Research|Quarterly|Press|Publications|MIT|Oxford|Cambridge|Springer|Wiley)",
        re.IGNORECASE,
    ),
    # Shorter citations with DOI (must be substantial, not just a bare URL)
    re.compile(r"[A-Z][a-z]+.*?doi[:\s/]", re.IGNORECASE),
    re.compile(r"[A-Z][a-z]+.*?https?://doi\.org/", re.IGNORECASE),
    # "Reference:" prefix pattern (as seen in CX decks)
    re.compile(r"^Reference:\s*.{30,}", re.IGNORECASE),
    # "Adapted from Author (Year)" pattern
    re.compile(r"^Adapted\s+from\s+[A-Z][a-z]+.*?\(\d{4}\)", re.IGNORECASE),
]

# Image attribution patterns
IMAGE_PATTERNS = [
    # Adobe Stock variants: "Adobe Stock 12345", "AdobeStock_12345", "Adobe Stock: 12345"
    re.compile(
        r"(?:Source:\s*)?(?:Image\s+(?:source|licensed|credit)[:\s]*)?(?:Adobe\s*Stock|AdobeStock)[_:\s#]*(\d+)",
        re.IGNORECASE,
    ),
    # "Image licensed through Adobe Stock: 12345"
    re.compile(
        r"Image\s+licensed\s+(?:through|from)\s+Adobe\s*Stock[_:\s#]*(\d+)",
        re.IGNORECASE,
    ),
    # Unsplash/Pexels/Getty/Shutterstock/iStock
    re.compile(
        r"(?:Source:\s*)?(?:Image\s+(?:source|credit|by)[:\s]*)?(.+?)\s+(?:on|via)\s+(Unsplash|Pexels|Getty|Shutterstock|iStock)",
        re.IGNORECASE,
    ),
    # Wikimedia Commons
    re.compile(
        r"(?:Source:\s*)?(?:Image\s+source:\s*)?(.*?),?\s*(?:Public\s+domain|CC\s+BY[-\s]SA).*?(?:via\s+)?Wikimedia\s+Commons",
        re.IGNORECASE,
    ),
    # "Diagram created by ..."
    re.compile(
        r"Diagram\s+created\s+by\s+(.+?)(?:\.\s*\(\d{4}\))?\.?\s*$",
        re.IGNORECASE,
    ),
    # "Video sourced via ..."
    re.compile(
        r"Video\s+sourced\s+(?:via|from)\s+(.+)",
        re.IGNORECASE,
    ),
    # Generic "Image source: ..." or "Photo credit: ..."
    re.compile(
        r"(?:Image|Photo)\s+(?:source|credit|by)[:\s]+(.+)",
        re.IGNORECASE,
    ),
]

# Source attribution pattern: "Source: XYZ"
SOURCE_PATTERN = re.compile(
    r"^Source:\s*(.+)",
    re.IGNORECASE,
)

# Noise to skip
NOISE_PATTERNS = [
    re.compile(r"(?i)^\s*$"),
    re.compile(r"(?i)^[\d\s]+$"),
    re.compile(r"(?i)click to (?:add|edit)"),
    re.compile(r"(?i)\[click"),
    re.compile(r"(?i)cricos"),
    re.compile(r"(?i)presentation\s+title"),
    re.compile(r"(?i)executive\s+education"),
]

# Footer-like text to skip
FOOTER_PATTERNS = [
    re.compile(r"(?i)^(?:UQ\s+)?Executive\s+Education"),
    re.compile(r"(?i)Negotiating\s+For\s+Success"),
    re.compile(r"(?i)Think\s+and\s+Act\s+Strategically"),
    re.compile(r"(?i)Leading\s+CX\s+Transformation"),
    re.compile(r"(?i)Climate\s+Finance\s+Essentials"),
]


def collect_references(source_prs: Presentation) -> dict:
    """
    Scan all slides in the source presentation for references and
    image attributions.

    Returns:
        {
            "academic": [
                {"text": "Full citation text...", "slide": 6},
                ...
            ],
            "images": [
                {"text": "Adobe Stock 12345", "slide": 14, "type": "adobe_stock"},
                {"text": "Maria Kovalets on Unsplash", "slide": 25, "type": "stock_photo"},
                ...
            ],
            "other_sources": [
                {"text": "Source: ASIC", "slide": 40},
                ...
            ],
        }
    """
    academic = []
    images = []
    other_sources = []

    # Track seen text to deduplicate
    seen_academic = set()
    seen_images = set()
    seen_other = set()

    for slide_idx, slide in enumerate(source_prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text or len(text) < 8:
                continue

            # Skip noise
            if any(p.search(text) for p in NOISE_PATTERNS):
                continue

            # Skip footer text
            if any(p.search(text) for p in FOOTER_PATTERNS):
                continue

            # Check each line in the text (some shapes have multiple lines)
            for line in text.split("\n"):
                line = line.strip()
                if not line or len(line) < 8:
                    continue
                if any(p.search(line) for p in NOISE_PATTERNS):
                    continue
                if any(p.search(line) for p in FOOTER_PATTERNS):
                    continue

                _classify_line(
                    line, slide_idx + 1,
                    academic, images, other_sources,
                    seen_academic, seen_images, seen_other,
                )

    return {
        "academic": academic,
        "images": images,
        "other_sources": other_sources,
    }


def _classify_line(
    line: str,
    slide_num: int,
    academic: list,
    images: list,
    other_sources: list,
    seen_academic: set,
    seen_images: set,
    seen_other: set,
):
    """Classify a single line of text as academic, image, or other source."""
    normalised = _normalise_for_dedup(line)

    # --- Image attributions (check first — more specific) ---
    for pattern in IMAGE_PATTERNS:
        if pattern.search(line):
            if normalised not in seen_images:
                seen_images.add(normalised)
                img_type = _classify_image_type(line)
                images.append({
                    "text": _clean_attribution(line),
                    "slide": slide_num,
                    "type": img_type,
                })
            return

    # --- "Source: ..." lines ---
    source_match = SOURCE_PATTERN.match(line)
    if source_match:
        source_text = source_match.group(1).strip()
        # Check if it's actually an image source (Adobe Stock in Source: line)
        if re.search(r"(?i)adobe\s*stock|unsplash|pexels|getty|shutterstock", source_text):
            if normalised not in seen_images:
                seen_images.add(normalised)
                images.append({
                    "text": _clean_attribution(line),
                    "slide": slide_num,
                    "type": _classify_image_type(line),
                })
        elif normalised not in seen_other:
            seen_other.add(normalised)
            other_sources.append({
                "text": _clean_attribution(line),
                "slide": slide_num,
            })
        return

    # --- Academic references ---
    for pattern in ACADEMIC_PATTERNS:
        if pattern.search(line):
            if normalised not in seen_academic:
                seen_academic.add(normalised)
                academic.append({
                    "text": _clean_citation(line),
                    "slide": slide_num,
                })
            return

    # --- Heuristic: looks like an APA citation ---
    if _looks_like_citation(line) and normalised not in seen_academic:
        seen_academic.add(normalised)
        academic.append({
            "text": _clean_citation(line),
            "slide": slide_num,
        })


def _looks_like_citation(text: str) -> bool:
    """
    Heuristic check for APA-ish citation patterns.
    Must have: author-like start + year + some length.
    Excludes profile URLs, biographical text, and YouTube links.
    """
    if len(text) < 40:
        return False

    lower = text.lower()

    # Reject common false positives
    if any(fp in lower for fp in [
        "youtube.com", "scholar.google", "business.uq.edu.au/profile",
        "researchers.uq.edu.au", "elected fellow", "appointed fellow",
        "from latin", "socialsciences.org.au",
    ]):
        return False

    # Reject if it's mostly a URL with no citation structure
    if re.match(r"^https?://", text) and "(" not in text:
        return False

    has_year = bool(re.search(r"\(\d{4}\)", text) or re.search(r",\s*\d{4}[,\.\)]", text))
    has_author = bool(re.match(r"[A-Z][a-z]+,?\s+[A-Z]", text))
    has_doi = bool(re.search(r"doi[:\s/]", text, re.IGNORECASE))
    has_journal = bool(re.search(r"(?i)journal|review|research|quarterly|press|publications", text))
    has_volume = bool(re.search(r"(?i)\bvol\b|\d+\(\d+\)", text))

    # Strong signals: year in parens + author-like start + journal/DOI
    if has_year and has_author and (has_journal or has_doi or has_volume):
        return True
    # DOI with some length
    if has_doi and len(text) > 60:
        return True
    # "Reference:" prefix already handled by ACADEMIC_PATTERNS
    return False


def _classify_image_type(text: str) -> str:
    """Classify the type of image attribution."""
    lower = text.lower()
    if "adobe" in lower or "stock" in lower:
        return "adobe_stock"
    if "unsplash" in lower:
        return "unsplash"
    if "pexels" in lower:
        return "pexels"
    if "getty" in lower or "istock" in lower:
        return "stock_photo"
    if "shutterstock" in lower:
        return "shutterstock"
    if "wikimedia" in lower:
        return "wikimedia"
    if "diagram" in lower:
        return "diagram"
    if "video" in lower:
        return "video"
    return "other"


def _clean_attribution(text: str) -> str:
    """Clean up an image attribution line."""
    text = text.strip()
    # Remove leading "Source: " if present (we'll add our own grouping)
    text = re.sub(r"^(?:Image\s+)?(?:source|credit)[:\s]*", "", text, flags=re.IGNORECASE).strip()
    if not text:
        return ""
    return text


def _clean_citation(text: str) -> str:
    """Clean up an academic citation line."""
    text = text.strip()
    # Remove leading "Reference:" if present
    text = re.sub(r"^Reference:\s*", "", text, flags=re.IGNORECASE).strip()
    return text


def _normalise_for_dedup(text: str) -> str:
    """
    Normalise text for deduplication.

    Uses first 80 chars after normalisation — catches the same citation
    appearing on multiple slides with slightly different trailing text
    (e.g. with/without DOI, with/without page range).
    """
    # Replace non-breaking spaces and other unicode whitespace
    text = text.replace("\xa0", " ").replace("\u200b", "")
    # Collapse whitespace
    text = re.sub(r"\s+", " ", text.lower().strip())
    # Strip trailing punctuation for dedup
    text = text.rstrip(".,;: ")
    # Truncate to first 80 chars — catches same citation with different endings
    return text[:80]


def format_references_text(refs: dict) -> str:
    """
    Format collected references into text suitable for a References slide.

    Groups: Academic References, then Image Credits, then Other Sources.
    Returns a single string with newline-separated entries.
    """
    parts = []

    if refs["academic"]:
        parts.append("References")
        for ref in refs["academic"]:
            parts.append(ref["text"])
        parts.append("")  # blank line separator

    if refs["images"]:
        parts.append("Image Credits")
        # Group by type
        adobe = [r for r in refs["images"] if r["type"] == "adobe_stock"]
        other_img = [r for r in refs["images"] if r["type"] != "adobe_stock"]

        if adobe:
            # Consolidate Adobe Stock IDs
            stock_ids = []
            for r in adobe:
                # Extract just the number
                match = re.search(r"(\d{5,})", r["text"])
                if match:
                    stock_ids.append(match.group(1))
                else:
                    stock_ids.append(r["text"])
            parts.append("Adobe Stock: " + ", ".join(stock_ids))

        for r in other_img:
            parts.append(r["text"])
        parts.append("")

    if refs["other_sources"]:
        parts.append("Other Sources")
        for ref in refs["other_sources"]:
            parts.append(ref["text"])

    return "\n".join(parts).strip()


def has_meaningful_references(refs: dict) -> bool:
    """Check if there are enough references to warrant a compiled slide."""
    total = len(refs["academic"]) + len(refs["images"]) + len(refs["other_sources"])
    return total >= 2
