"""
Table of Contents Generator.

Scans converted slides for Section Divider layouts and builds a
"Contents" slide using the Contents 2 template layout (index 4).

The Contents 2 layout has a table placeholder (idx 12) that receives a
two-column table: numbered section entries (01, 02, ...) paired with
section titles.  This mirrors the pattern found in human-branded decks
such as Climate Finance Essentials.

The generated slide is inserted early in the deck — after Cover and
Acknowledgement of Country but before the first content slide.

Design principles:
    - Only generate if 3+ section dividers are found
    - Use auto-numbering (01, 02, ...) regardless of source numbering
    - NEVER set font names, sizes, or colours — inherit from template
    - Detect and skip existing Contents/Agenda slides in the source
"""

import logging
import re

from pptx.presentation import Presentation

logger = logging.getLogger("uqslide.toc")


# Layout names that indicate a section divider in the converted output
SECTION_DIVIDER_LAYOUT_NAMES = [
    "section divider",
    "section header",
]

# Layout names / title text that indicate an existing contents slide
# (should not be counted as a section divider)
CONTENTS_INDICATORS = [
    "contents",
    "table of contents",
    "agenda",
    "course overview",
    "outline",
]

# Template layout index for Contents 2
CONTENTS_2_LAYOUT_INDEX = 4
# Placeholder indices in Contents 2
PH_TITLE = 0
PH_TABLE = 12
PH_FOOTER = 17
PH_SLIDE_NUM = 18

# Minimum sections required to generate a ToC
MIN_SECTIONS = 3


def collect_sections(prs: Presentation) -> list[dict]:
    """
    Scan all slides in the converted presentation for Section Divider
    layouts and extract their titles and section numbers.

    Returns a list of dicts:
        [
            {"title": "Introduction to climate finance", "section_num": "02", "slide_index": 5},
            ...
        ]
    """
    sections = []

    for slide_idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name.lower()

        # Only interested in section divider layouts
        if not any(d in layout_name for d in SECTION_DIVIDER_LAYOUT_NAMES):
            continue

        # Skip if this looks like a contents/agenda slide (shouldn't happen
        # with Section Divider layout, but be safe)
        title_text = _get_slide_title(slide)
        if title_text and any(kw in title_text.lower() for kw in CONTENTS_INDICATORS):
            continue

        section_num = ""
        section_title = title_text or ""

        # Try to extract the section number from placeholder idx 11
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        if 11 in placeholders:
            num_text = placeholders[11].text.strip()
            if num_text:
                section_num = num_text

        # Get title from placeholder idx 0 if we didn't get it above
        if not section_title and 0 in placeholders:
            section_title = placeholders[0].text.strip()

        # Robustness: if section_num looks like a title (long text) and
        # title looks like a label (short), swap them. This handles
        # human-branded decks where placeholders may be used differently.
        if section_num and section_title:
            num_looks_like_label = _looks_like_section_label(section_num)
            title_looks_like_label = _looks_like_section_label(section_title)
            if not num_looks_like_label and title_looks_like_label:
                # Swap: the "num" field has the real title
                section_num, section_title = section_title, section_num

        if section_title:
            sections.append({
                "title": section_title,
                "section_num": section_num,
                "slide_index": slide_idx,
            })

    return sections


def should_generate_toc(sections: list[dict]) -> bool:
    """Check if there are enough sections to warrant a ToC slide."""
    return len(sections) >= MIN_SECTIONS


def build_toc_content(sections: list[dict]) -> list[tuple[str, str]]:
    """
    Build the table rows for the ToC slide.

    Always uses auto-numbering (01, 02, ...) for consistency,
    regardless of what the source section numbers were.

    Returns list of (number_str, title_str) tuples.
    """
    rows = []
    for i, section in enumerate(sections, start=1):
        num = f"{i:02d}"
        title = section["title"]
        rows.append((num, title))
    return rows


def _get_slide_title(slide) -> str:
    """Get the title text from a slide (placeholder idx 0 or shape named Title)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if (hasattr(shape, 'is_placeholder') and shape.is_placeholder
                and shape.placeholder_format.idx == 0):
            text = shape.text_frame.text.strip()
            if text:
                return text

    # Fallback: shape named "Title"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "title" in shape.name.lower():
            text = shape.text_frame.text.strip()
            if text:
                return text

    return ""


def _looks_like_section_label(text: str) -> bool:
    """
    Check if text looks like a section number/label rather than a title.

    Section labels: "01", "Session 1:", "Block 3", "Module 2", etc.
    Titles: "Introduction to climate finance", "Strategic Challenges", etc.
    """
    text = text.strip()
    # Zero-padded numbers
    if re.match(r"^0\d{1,2}$", text):
        return True
    # "Session N:", "Block N", etc.
    if re.match(r"(?i)^(session|block|module|part|week|day|unit|topic|section)\s*\d", text):
        return True
    # Very short (< 15 chars) and contains a digit
    if len(text) < 15 and re.search(r"\d", text):
        return True
    return False


def _has_existing_toc(prs: Presentation) -> bool:
    """
    Check if the presentation already has a Contents/Agenda slide.

    Looks at the first 5 slides (after cover and AoC) for:
    - Slides using a "Contents" layout
    - Slides with "Contents" or "Agenda" in the title
    """
    for slide_idx in range(min(5, len(prs.slides))):
        slide = prs.slides[slide_idx]

        # Check layout name
        layout_name = slide.slide_layout.name.lower()
        if any(kw in layout_name for kw in ["contents", "agenda"]):
            return True

        # Check title text
        title = _get_slide_title(slide)
        if title and any(kw in title.lower() for kw in CONTENTS_INDICATORS):
            return True

    return False
