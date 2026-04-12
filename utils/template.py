"""
Template utilities for creating output slides from the UQ template.
"""

import os
from pptx import Presentation

# Template lives alongside the app
TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates")
TEMPLATE_FILENAME = "preferred_template.pptx"


def open_template() -> Presentation:
    """Open a fresh copy of the UQ template presentation."""
    path = os.path.join(TEMPLATE_DIR, TEMPLATE_FILENAME)
    return Presentation(path)


def add_slide_from_layout(prs: Presentation, layout_index: int):
    """
    Add a new slide using the specified layout index from the template.
    Also activates footer and slide number placeholders from the layout
    (python-pptx doesn't inherit these automatically).
    Returns the new slide object.
    """
    layout = prs.slide_masters[0].slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    # Activate footer/slide number placeholders from layout
    _activate_layout_placeholders(slide, layout, ph_types=["ftr", "sldNum"])

    return slide


def _activate_layout_placeholders(slide, layout, ph_types):
    """
    Copy placeholder shapes from the layout XML to the slide XML.
    This makes inherited footer/slide number placeholders editable.

    Args:
        slide: the new slide
        layout: the slide layout
        ph_types: list of placeholder type strings to activate
                  ("ftr" = footer, "sldNum" = slide number)
    """
    from copy import deepcopy
    from lxml import etree

    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Find existing placeholder types on the slide
    existing_types = set()
    for sp in slide.shapes._spTree.findall('.//p:sp', nsmap):
        ph = sp.find('.//p:nvSpPr/p:nvPr/p:ph', nsmap)
        if ph is not None:
            existing_types.add(ph.get('type'))

    # Copy missing placeholders from layout
    for sp in layout.placeholders._element.getparent().findall('.//p:sp', nsmap):
        ph = sp.find('.//p:nvSpPr/p:nvPr/p:ph', nsmap)
        if ph is not None:
            ph_type = ph.get('type')
            if ph_type in ph_types and ph_type not in existing_types:
                slide.shapes._spTree.append(deepcopy(sp))


def delete_slide(prs: Presentation, slide_index: int):
    """Delete a slide by index using XML manipulation."""
    rId = prs.slides._sldIdLst[slide_index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(sldId)


def delete_all_original_slides(prs: Presentation, num_new_slides: int = 1):
    """
    Delete all the original template slides, keeping only the newly added ones.
    The new slides are always appended at the end, so we delete from the
    beginning (in reverse order to preserve indices).
    """
    num_original = len(prs.slides) - num_new_slides
    for i in range(num_original - 1, -1, -1):
        delete_slide(prs, i)


def move_slide_to_position(prs: Presentation, from_index: int, to_index: int):
    """
    Move a slide from one position to another by reordering the sldIdLst XML.

    Args:
        prs: Presentation object
        from_index: Current 0-based index of the slide to move
        to_index: Target 0-based index where the slide should end up
    """
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[from_index]
    sldIdLst.remove(sldId)
    sldIdLst.insert(to_index, sldId)
